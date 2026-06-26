"""Tests for template-first board PPTX export."""

from __future__ import annotations

from app.services.reporting.export.pptx_template_export import (
    TemplateCommentaryUpdate,
    _apply_key_takeaway_bullets,
    _find_key_takeaways_slots,
    _find_commentary_column_shapes,
    _find_key_takeaways_bullet_shapes,
    _find_risk_description_shapes,
    _is_boilerplate_shape,
    _is_exec_nav_text,
    _month_label,
    _normalize_bullet_line,
    _period_replacements,
    _prior_month,
    _quarter_label,
    _set_shape_text,
    _split_narrative_into_bullets,
    _truncate_to_char_limit,
    _ytd_range_label,
    apply_template_commentary,
    map_template_slides_to_slide_keys,
    repair_executive_nav_links,
    resolve_board_pptx_template,
    roll_template_period_labels,
)
from app.services.reporting.export.schemas import ExportValidationSummary, ReportingBundle
from app.services.dashboard.schemas import ExecutiveFlowResponse


def _minimal_bundle(*, as_of: str = "2026-06") -> ReportingBundle:
    org = "00000000-0000-0000-0000-000000000001"
    return ReportingBundle(
        organization_id=org,
        organization_name="Test Co",
        scenario="Combined",
        period_label=_month_label(as_of),
        as_of_period=as_of,
        start_period="2026-01",
        end_period="2026-12",
        currency="USD",
        executive_flow=ExecutiveFlowResponse(
            organization_id=org,
            scenario="Combined",
            start_period="2026-01",
            end_period="2026-12",
        ),
        validation=ExportValidationSummary(status="pass", failed_count=0, warning_count=0, passed_count=0, checks=[]),
    )


def test_period_label_helpers():
    assert _month_label("2026-06") == "June 2026"
    assert _prior_month("2026-06") == "2026-05"
    assert _quarter_label("2026-06") == "Q2 2026"
    assert _ytd_range_label("2026-06") == "Jan-Jun 2026"


def test_period_replacements_roll_may_to_june():
    bundle = _minimal_bundle(as_of="2026-06")
    patterns = _period_replacements(bundle)
    sample = "Board Review — May 2026 | Q2 2026 | Jan-May 2026"
    rolled = sample
    for pattern, repl in patterns:
        rolled = pattern.sub(repl, rolled)
    assert "June 2026" in rolled
    assert "May 2026" not in rolled
    assert "Jan-Jun 2026" in rolled


def test_resolve_template_returns_none_when_missing(monkeypatch, tmp_path):
    monkeypatch.setattr(
        "app.services.reporting.export.pptx_template_export._TEMPLATE_SEARCH_GLOBS",
        ((tmp_path, "*.pptx"),),
    )
    assert resolve_board_pptx_template() is None


def test_exec_nav_text_detection():
    assert _is_exec_nav_text("View waterfall →")
    assert _is_exec_nav_text("View P&L →")
    assert not _is_exec_nav_text("June 2026 close")


def test_roll_preserves_exec_nav_paragraph():
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    tf = box.text_frame
    tf.paragraphs[0].text = "Board Review — May 2026"
    tf.add_paragraph().text = "View waterfall →"

    bundle = _minimal_bundle(as_of="2026-06")
    roll_template_period_labels(prs, bundle)

    assert "June 2026" in box.text
    assert "View waterfall" in box.text
    assert "May 2026" not in box.text


def test_repair_executive_nav_links_binds_internal_jumps():
    from pptx import Presentation
    from pptx.enum.action import PP_ACTION
    from pptx.util import Inches

    prs = Presentation()
    exec_slide = prs.slides.add_slide(prs.slide_layouts[6])
    exec_box = exec_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1.5))
    exec_tf = exec_box.text_frame
    exec_tf.paragraphs[0].text = "Executive Operating Summary"
    exec_tf.add_paragraph().text = "View waterfall →"
    exec_tf.add_paragraph().text = "View driver tree →"
    exec_tf.add_paragraph().text = "View P&L →"
    exec_tf.add_paragraph().text = "View cash →"

    arr_slide = prs.slides.add_slide(prs.slide_layouts[6])
    arr_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6), Inches(1)).text = "ARR / MRR Waterfall"
    driver_slide = prs.slides.add_slide(prs.slide_layouts[6])
    driver_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6), Inches(1)).text = "Net New ARR driver decomposition"
    pl_slide = prs.slides.add_slide(prs.slide_layouts[6])
    pl_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6), Inches(1)).text = "Department Spend - P&L"
    cash_slide = prs.slides.add_slide(prs.slide_layouts[6])
    cash_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6), Inches(1)).text = "Cash Forecast"

    repaired = repair_executive_nav_links(prs)
    assert repaired == 4

    for label in ("View waterfall", "View driver", "View P&L", "View cash"):
        found = False
        for paragraph in exec_tf.paragraphs:
            for run in paragraph.runs:
                if label.lower() in (run.text or "").lower():
                    r_pr = run._r.rPr
                    assert r_pr is not None and r_pr.hlinkClick is not None
                    assert r_pr.hlinkClick.action == "ppaction://hlinksldjump"
                    assert r_pr.hlinkClick.rId
                    found = True
        assert found, label

    assert exec_box.click_action.action == PP_ACTION.NONE


def test_map_template_slides_to_slide_keys():
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slides_text = [
        "Executive Summary — Where we stand",
        "ARR Analysis / MRR Waterfall",
        "P&L Review — Income Statement",
        "Cash & Liquidity — Cash Forecast",
        "GTM & Marketing — Channel efficiency",
        "Workforce & Headcount",
        "Risks & Opportunities",
    ]
    for title in slides_text:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1)).text = title

    mapping = map_template_slides_to_slide_keys(prs)
    assert mapping["executive_summary"] == 1
    assert mapping["arr_waterfall"] == 2
    assert mapping["gaap_revenue"] == 3
    assert mapping["cash_forecast"] == 4
    assert mapping["gtm_performance"] == 5
    assert mapping["headcount"] == 6
    assert mapping["risks_opportunities"] == 7


def test_set_shape_text_preserves_multiline_bullets():
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    _set_shape_text(box, "Key Takeaways\n• Revenue $7.41M vs budget\n• ARR $86.10M")
    assert box.text_frame.paragraphs[0].text == "Key Takeaways"
    assert "Revenue" in box.text_frame.paragraphs[1].text
    assert "ARR" in box.text_frame.paragraphs[2].text


def test_is_boilerplate_shape_detects_footer():
    assert _is_boilerplate_shape("SMPL · Board Operating Review · Q2 2026 · CONFIDENTIAL")
    assert _is_boilerplate_shape("2/11")
    assert not _is_boilerplate_shape("• Revenue $7.41M vs budget")


def test_apply_key_takeaway_bullets_fills_template_slots_in_order():
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(0.4)).text = "Key Takeaways"
    boxes = []
    for row in range(3):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0 + row * 0.5), Inches(8), Inches(0.45))
        box.text = f"• Old template bullet {row + 1} about May revenue"
        boxes.append(box)
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(8), Inches(0.4))
    footer.text = "SMPL · Board Operating Review · Q2 2026 · CONFIDENTIAL"

    slots = _find_key_takeaways_slots(slide)
    assert len(slots) == 3

    bullets = ["• First provided bullet", "• Second provided bullet", "• Third provided bullet"]
    assert _apply_key_takeaway_bullets(slide, bullets)

    assert "May" not in boxes[0].text
    assert "First provided" in boxes[0].text
    assert "Second provided" in boxes[1].text
    assert "Third provided" in boxes[2].text
    assert footer.text == "SMPL · Board Operating Review · Q2 2026 · CONFIDENTIAL"


def test_apply_key_takeaway_bullets_clears_extra_template_slots():
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(0.4)).text = "Key Takeaways"
    boxes = []
    for row in range(3):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0 + row * 0.5), Inches(8), Inches(0.45))
        box.text = f"• Old bullet {row + 1}"
        boxes.append(box)

    assert _apply_key_takeaway_bullets(slide, ["• Only one new bullet"])
    assert "Only one new" in boxes[0].text
    assert boxes[1].text.strip() == ""
    assert boxes[2].text.strip() == ""


def test_normalize_bullet_line_strips_duplicate_prefix():
    assert _normalize_bullet_line("• Revenue $7.35M") == "• Revenue $7.35M"
    assert _normalize_bullet_line("Revenue $7.35M") == "• Revenue $7.35M"


def test_truncate_to_char_limit_preserves_word_boundary():
    long_text = "• Ending ARR $86.10M vs budget $84.50M; net new ARR $1.60M driven by expansion and new logos in June."
    trimmed = _truncate_to_char_limit(long_text, 90)
    assert len(trimmed) <= 90
    assert trimmed.endswith("…")


def test_split_narrative_into_bullets_chunks_long_prose():
    prose = (
        "Operating cash flow was negative in June due to payroll timing and annual prepayments. "
        "Collections remained on plan. Financing inflows offset the bridge. "
        "H2 outlook assumes normalized working capital."
    )
    chunks = _split_narrative_into_bullets(prose, max_count=4, char_limit=90)
    assert len(chunks) >= 2
    assert all(len(c) <= 90 for c in chunks)


def test_apply_key_takeaway_bullets_finds_empty_slots_after_clear():
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(0.4)).text = "Key Takeaways"
    boxes = []
    for row in range(4):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0 + row * 0.5), Inches(8), Inches(0.45))
        box.text = f"• Old bullet {row + 1}"
        boxes.append(box)

    assert _apply_key_takeaway_bullets(slide, ["• Only one new bullet"], slide_key="cash_flow_statement")
    assert "Only one new" in boxes[0].text
    assert boxes[1].text.strip() == ""
    slots = _find_key_takeaways_slots(slide)
    assert len(slots) == 4


def test_apply_template_commentary_skips_footer_only_slides():
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(8), Inches(0.4))
    footer.text = "SMPL · Board Operating Review · Q2 2026 · CONFIDENTIAL"
    table_cell = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(4), Inches(0.4))
    table_cell.text = "1.3x, 5.8% WR — reallocate to Partner"

    applied = apply_template_commentary(
        prs,
        {
            1: TemplateCommentaryUpdate(
                slide_key="executive_summary",
                text="• Fresh June bullet",
                bullets=("• Fresh June bullet",),
            )
        },
    )
    assert applied == 0
    assert "Fresh June" not in footer.text
    assert "Fresh June" not in table_cell.text


def test_map_template_slides_includes_gtm_and_outlook_on_gold_deck():
    template = resolve_board_pptx_template()
    if template is None:
        return
    from pptx import Presentation

    mapping = map_template_slides_to_slide_keys(Presentation(str(template)))
    assert mapping.get("executive_summary") == 2
    assert mapping.get("gtm_performance") == 7
    assert mapping.get("gtm_funnel") == 8
    assert mapping.get("risks_opportunities") == 9
    assert mapping.get("financial_outlook") == 10
    assert mapping.get("board_actions") == 11
    assert "headcount" not in mapping
