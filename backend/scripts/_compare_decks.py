"""Detailed slide-by-slide diff: v3 reference vs recreated deck."""
from __future__ import annotations

import sys
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Emu

REF = Path(
    r"C:\Users\mattj\.cursor\projects\empty-window\saas-financial-intelligence\backend\tmp\mda_deck_2026-06 v3.pptx"
)
NEW = Path(
    r"C:\Users\mattj\.cursor\projects\empty-window\saas-financial-intelligence\backend\tmp\mda_deck_2026-06_recreated.pptx"
)


@dataclass
class ShapeInfo:
    kind: str
    text: str = ""
    chart: str = ""
    table_rows: list[list[str]] = field(default_factory=list)
    box: tuple[int, int, int, int] = (0, 0, 0, 0)  # x,y,w,h in EMU


def _emu(v) -> int:
    return int(v) if v is not None else 0


def _norm(s: str) -> str:
    return " ".join(s.split()).strip()


def extract_slide(slide) -> list[ShapeInfo]:
    out: list[ShapeInfo] = []
    for shape in slide.shapes:
        x, y, w, h = _emu(shape.left), _emu(shape.top), _emu(shape.width), _emu(shape.height)
        if shape.has_table:
            rows: list[list[str]] = []
            for r in shape.table.rows:
                rows.append([_norm(c.text) for c in r.cells])
            out.append(ShapeInfo("table", table_rows=rows, box=(x, y, w, h)))
        elif shape.has_chart:
            ch = shape.chart
            series = []
            for s in ch.series:
                pts = []
                try:
                    pts = [round(float(v), 3) for v in s.values]
                except Exception:
                    pts = list(s.values)
                series.append({"name": s.name, "values": pts})
            cats = []
            try:
                cats = [str(c) for c in ch.plots[0].categories]
            except Exception:
                pass
            chart_desc = (
                f"{XL_CHART_TYPE(ch.chart_type).name}|cats={cats}|series={series}"
            )
            out.append(ShapeInfo("chart", chart=chart_desc, box=(x, y, w, h)))
        elif hasattr(shape, "text") and shape.text.strip():
            out.append(ShapeInfo("text", text=_norm(shape.text), box=(x, y, w, h)))
        elif shape.shape_type is not None:
            out.append(ShapeInfo(f"shape:{shape.shape_type}", box=(x, y, w, h)))
    out.sort(key=lambda s: (s.box[1], s.box[0]))
    return out


def summarize(slide_infos: list[ShapeInfo]) -> dict:
    texts = [s.text for s in slide_infos if s.kind == "text"]
    charts = [s.chart for s in slide_infos if s.kind == "chart"]
    tables = [s.table_rows for s in slide_infos if s.kind == "table"]
    return {"texts": texts, "charts": charts, "tables": tables, "n_shapes": len(slide_infos)}


def diff_lists(a: list, b: list, label: str) -> list[str]:
    diffs: list[str] = []
    if len(a) != len(b):
        diffs.append(f"  {label} count: ref={len(a)} new={len(b)}")
    for i, (x, y) in enumerate(zip(a, b)):
        if x != y:
            diffs.append(f"  {label}[{i}] ref: {str(x)[:200]}")
            diffs.append(f"  {label}[{i}] new: {str(y)[:200]}")
    if len(a) > len(b):
        for x in a[len(b) :]:
            diffs.append(f"  {label} only in ref: {str(x)[:200]}")
    elif len(b) > len(a):
        for y in b[len(a) :]:
            diffs.append(f"  {label} only in new: {str(y)[:200]}")
    return diffs


def main() -> None:
    sys.stdout.reconfigure(encoding="utf-8")
    ref_p = Presentation(str(REF))
    new_p = Presentation(str(NEW))
    print(f"Reference: {REF.name} ({REF.stat().st_size:,} bytes, {len(ref_p.slides)} slides)")
    print(f"Recreated: {NEW.name} ({NEW.stat().st_size:,} bytes, {len(new_p.slides)} slides)")
    print()

    n = max(len(ref_p.slides), len(new_p.slides))
    slide_diffs = 0
    for i in range(n):
        ref_slide = ref_p.slides[i] if i < len(ref_p.slides) else None
        new_slide = new_p.slides[i] if i < len(new_p.slides) else None
        ref_info = extract_slide(ref_slide) if ref_slide else []
        new_info = extract_slide(new_slide) if new_slide else []
        ref_sum = summarize(ref_info)
        new_sum = summarize(new_info)

        diffs: list[str] = []
        diffs.extend(diff_lists(ref_sum["texts"], new_sum["texts"], "text"))
        diffs.extend(diff_lists(ref_sum["charts"], new_sum["charts"], "chart"))
        diffs.extend(diff_lists(ref_sum["tables"], new_sum["tables"], "table"))

        if ref_sum["n_shapes"] != new_sum["n_shapes"]:
            diffs.insert(0, f"  shape count: ref={ref_sum['n_shapes']} new={new_sum['n_shapes']}")

        if diffs:
            slide_diffs += 1
            print(f"--- SLIDE {i + 1} — DIFFERENCES ---")
            for d in diffs:
                print(d)
            print()

    if slide_diffs == 0:
        print("No differences detected.")
    else:
        print(f"Summary: {slide_diffs}/{n} slides differ.")


if __name__ == "__main__":
    main()
