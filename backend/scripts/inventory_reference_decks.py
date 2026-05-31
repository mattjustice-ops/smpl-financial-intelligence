"""Inventory reference-decks PPTX/XLSX for docs/reference-decks/INDEX.md."""
from __future__ import annotations

import sys
from pathlib import Path

BASE = Path(__file__).resolve().parents[2] / "docs" / "reference-decks"
OUT = BASE / "_inventory.txt"


def main() -> int:
    lines: list[str] = []

    pptx = BASE / "ClarityFP_Board_Review_May2026_v5.pptx"
    if pptx.exists():
        from pptx import Presentation

        prs = Presentation(str(pptx))
        lines.append(f"PPTX: {pptx.name} ({len(prs.slides)} slides)")
        for i, slide in enumerate(prs.slides, 1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip().replace("\n", " | ")[:140])
            lines.append(f"  {i:2d}: {texts[0] if texts else '(no text)'}")
    else:
        lines.append(f"MISSING: {pptx.name}")

    try:
        import openpyxl
    except ImportError:
        lines.append("openpyxl not installed")
        openpyxl = None  # type: ignore

    for name in sorted(BASE.glob("*.xlsx")):
        if not openpyxl:
            break
        wb = openpyxl.load_workbook(name, read_only=True, data_only=True)
        lines.append(f"\nXLSX: {name.name}")
        lines.append(f"  sheets: {', '.join(wb.sheetnames)}")
        for sn in wb.sheetnames[:6]:
            ws = wb[sn]
            preview: list[str] = []
            for row in ws.iter_rows(max_row=5, max_col=10, values_only=True):
                cells = [str(c)[:20] if c is not None else "" for c in row]
                if any(cells):
                    preview.append(" | ".join(cells))
            if preview:
                lines.append(f"  [{sn}]")
                lines.extend(f"    {p[:180]}" for p in preview[:3])
        wb.close()

    OUT.write_text("\n".join(lines), encoding="utf-8")
    print(f"Wrote {OUT}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
