"""Fill an existing board PPTX template — roll period labels and inject Claude commentary.

The reference deck (ClarityFP / SMPL Board Review) keeps layout, charts, and visuals.
We update month/quarter labels and replace narrative text blocks per slide.
"""

from __future__ import annotations

import calendar
import io
import json
import logging
import re
from dataclasses import dataclass
from datetime import date
from pathlib import Path
from typing import Any, Iterator

from app.core.config import _BACKEND_ROOT, get_settings
from app.services.commentary.llm_factory import build_commentary_llm_client
from app.services.reporting.export.board_commentary_service import (
    build_all_slide_commentary,
    copilot_context_blob,
)
from app.services.reporting.export.export_sheet_registry import requirements_prompt_block
from app.services.reporting.export.schemas import ReportingBundle

logger = logging.getLogger(__name__)

_REPO_ROOT = _BACKEND_ROOT.parent

# Search order for the gold-standard deck (user-provided template).
_CANONICAL_TEMPLATE_NAME = "SMPL_Board_Review_Template.pptx"
_CANONICAL_TEMPLATE = _BACKEND_ROOT / "templates" / "board" / _CANONICAL_TEMPLATE_NAME
_TEMPLATE_SEARCH_GLOBS: tuple[tuple[Path, str], ...] = (
    (_BACKEND_ROOT / "templates" / "board", "SMPL_Board_Review*.pptx"),
    (_REPO_ROOT / "docs" / "reference-decks", "SMPL_Board_Review*.pptx"),
    (_REPO_ROOT / "frontend" / "public" / "board" / "exports", "SMPL_Board_Review*.pptx"),
)

_ONEDRIVE_TEMPLATE = Path.home() / "OneDrive" / "SMPL_Board_Review_Q2_2026.pptx"


@dataclass(frozen=True)
class TemplateSlideOutline:
    index: int
    title: str
    preview: str


@dataclass(frozen=True)
class TemplateCommentaryUpdate:
    slide_key: str
    text: str
    bullets: tuple[str, ...] = ()


# How commentary text is placed into each mapped slide's template zones.
_SLIDE_COMMENTARY_ZONE: dict[str, str] = {
    "executive_summary": "takeaway_bullets",
    "arr_waterfall": "takeaway_bullets",
    "gaap_revenue": "takeaway_bullets",
    "cash_forecast": "takeaway_bullets",
    "cash_flow_statement": "takeaway_bullets",
    "gtm_performance": "commentary_column",
    "gtm_funnel": "narrative_blocks",
    "headcount": "takeaway_bullets",
    "risks_opportunities": "risk_descriptions",
    "financial_outlook": "narrative_blocks",
    "board_actions": "narrative_blocks",
}

# Template box capacity (chars) — derived from gold deck geometry; export layer enforces hard caps.
_SLIDE_SLOT_CHAR_LIMIT: dict[str, int] = {
    "executive_summary": 100,
    "arr_waterfall": 90,
    "gaap_revenue": 85,
    "cash_forecast": 90,
    "cash_flow_statement": 90,
    "gtm_performance": 38,
    "gtm_funnel": 115,
    "risks_opportunities": 135,
    "financial_outlook": 160,
    "board_actions": 165,
}

_KEY_TAKEAWAY_SLOT_GAP_EMU = 28000
_TEXT_FRAME_MARGIN_EMU = 45000
_GTM_CHANNEL_COL_LEFT = 411480
_GTM_ROW_TOP_TOLERANCE = 200000


def resolve_board_pptx_template() -> Path | None:
    settings = get_settings()
    explicit = getattr(settings, "board_pptx_template", None)
    if explicit:
        path = Path(explicit)
        if path.is_file():
            logger.info("Board PPTX template (BOARD_PPTX_TEMPLATE): %s", path)
            return path

    if _CANONICAL_TEMPLATE.is_file():
        logger.info("Board PPTX template (canonical): %s", _CANONICAL_TEMPLATE)
        return _CANONICAL_TEMPLATE

    if _ONEDRIVE_TEMPLATE.is_file():
        logger.info("Board PPTX template (OneDrive dev): %s", _ONEDRIVE_TEMPLATE)
        return _ONEDRIVE_TEMPLATE

    best: Path | None = None
    best_slides = 0
    from pptx import Presentation

    for folder, pattern in _TEMPLATE_SEARCH_GLOBS:
        if not folder.is_dir():
            continue
        for candidate in sorted(folder.glob(pattern), key=lambda p: p.stat().st_mtime, reverse=True):
            try:
                slide_count = len(Presentation(str(candidate)).slides)
            except Exception:
                continue
            if slide_count >= 10 and slide_count >= best_slides:
                best = candidate
                best_slides = slide_count
    if best is not None:
        logger.info("Board PPTX template (glob): %s (%s slides)", best, best_slides)
    return best


def _month_label(period: str) -> str:
    dt = date(int(period[:4]), int(period[5:7]), 1)
    return dt.strftime("%B %Y")


def _prior_month(period: str) -> str:
    year = int(period[:4])
    month = int(period[5:7])
    if month == 1:
        return f"{year - 1}-12"
    return f"{year}-{month - 1:02d}"


def _quarter_label(period: str) -> str:
    month = int(period[5:7])
    quarter = (month - 1) // 3 + 1
    return f"Q{quarter} {period[:4]}"


def _ytd_range_label(period: str) -> str:
    month = int(period[5:7])
    short = calendar.month_abbr[month]
    return f"Jan-{short} {period[:4]}"


def _period_replacements(bundle: ReportingBundle) -> list[tuple[re.Pattern[str], str]]:
    as_of = bundle.as_of_period
    close = _month_label(as_of)
    prior = _month_label(_prior_month(as_of))
    quarter = _quarter_label(as_of)
    ytd = _ytd_range_label(as_of)
    year = as_of[:4]

    patterns: list[tuple[str, str]] = [
        (
            r"\b(?:Jan(?:uary)?|Feb(?:ruary)?|Mar(?:ch)?|Apr(?:il)?|May|Jun(?:e)?|Jul(?:y)?|Aug(?:ust)?|Sep(?:tember)?|Oct(?:ober)?|Nov(?:ember)?|Dec(?:ember)?)\s+20\d{2}\b",
            close,
        ),
        (r"\bQ[1-4]\s+20\d{2}\b", quarter),
        (r"\bJan-[A-Za-z]{3}\s+20\d{2}\b", ytd),
        (r"\b20\d{2}\s+Q[1-4]\b", f"{year} {quarter.split()[0]}"),
        (r"\bClose:\s*[A-Za-z]+\s+20\d{2}\b", f"Close: {close}"),
        (r"\bClose\s+[A-Za-z]+\s+20\d{2}\b", f"Close {close}"),
    ]
    out: list[tuple[re.Pattern[str], str]] = []
    for pattern, repl in patterns:
        out.append((re.compile(pattern, re.I), repl))
    out.append((re.compile(re.escape(prior), re.I), close))
    return out


def _iter_shapes(shapes) -> Iterator[Any]:
    """Walk slide shapes including grouped children."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        MSO_SHAPE_TYPE = None  # type: ignore[assignment,misc]

    for shape in shapes:
        try:
            if MSO_SHAPE_TYPE is not None and getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
                yield from _iter_shapes(shape.shapes)
                continue
        except Exception:
            pass
        yield shape


def _shape_text(shape) -> str:
    try:
        if not hasattr(shape, "text"):
            return ""
        return (shape.text or "").strip()
    except Exception:
        return ""


def _slide_text_blob(slide) -> str:
    parts = [_shape_text(s) for s in _iter_shapes(slide.shapes) if _shape_text(s)]
    return " ".join(parts)


_EXEC_NAV_LINKS: tuple[tuple[re.Pattern[str], tuple[str, ...]], ...] = (
    (re.compile(r"view\s+waterfall", re.I), ("arr / mrr waterfall", "mrr waterfall", "arr waterfall", "arr roll")),
    (re.compile(r"view\s+driver", re.I), ("driver", "net new arr", "why it happened", "decomposition")),
    (re.compile(r"view\s+p\s*&?\s*l", re.I), ("department spend", "management p&l", "p&l", "income statement")),
    (re.compile(r"view\s+cash", re.I), ("cash forecast", "cash flow", "liquidity", "cash bridge")),
)


def _is_exec_nav_text(text: str) -> bool:
    lowered = text.lower()
    return any(
        token in lowered
        for token in ("view waterfall", "view driver", "view p&l", "view pl", "view cash")
    )


def _paragraph_has_slide_link(paragraph) -> bool:
    for run in paragraph.runs:
        r_pr = run._r.rPr
        if r_pr is None or r_pr.hlinkClick is None:
            continue
        hlink = r_pr.hlinkClick
        action = (hlink.action or "").lower()
        if "hlinksldjump" in action or hlink.rId:
            return True
    return False


def _shape_has_slide_link(shape) -> bool:
    try:
        from pptx.enum.action import PP_ACTION

        if shape.click_action.action != PP_ACTION.NONE:
            return True
    except Exception:
        pass
    if not hasattr(shape, "text_frame") or shape.text_frame is None:
        return False
    return any(_paragraph_has_slide_link(para) for para in shape.text_frame.paragraphs)


def _is_protected_nav_shape(shape) -> bool:
    text = _shape_text(shape)
    if _is_exec_nav_text(text):
        return True
    return _shape_has_slide_link(shape)


def _slide_keyword_index(prs) -> list[tuple[int, str]]:
    return [(idx, _slide_text_blob(slide).lower()) for idx, slide in enumerate(prs.slides)]


def _find_slide_by_keywords(
    prs,
    keywords: tuple[str, ...],
    *,
    exclude_indices: set[int] | None = None,
) -> int | None:
    exclude = exclude_indices or set()
    best_idx: int | None = None
    best_score = 0
    for idx, blob in _slide_keyword_index(prs):
        if idx in exclude:
            continue
        score = sum(1 for keyword in keywords if keyword in blob)
        if score > best_score:
            best_score = score
            best_idx = idx
    return best_idx if best_score > 0 else None


def _find_executive_summary_slide(prs) -> tuple[Any | None, int | None]:
    best_slide: Any | None = None
    best_idx: int | None = None
    best_score = 0
    for idx, slide in enumerate(prs.slides):
        blob = _slide_text_blob(slide).lower()
        link_hits = sum(1 for pattern, _ in _EXEC_NAV_LINKS if pattern.search(blob))
        score = link_hits
        if "executive" in blob or "operating summary" in blob:
            score += 2
        if "where we stand" in blob:
            score += 1
        if score > best_score:
            best_score = score
            best_slide = slide
            best_idx = idx
    if best_score >= 2:
        return best_slide, best_idx
    return None, None


def _set_run_target_slide(run, target_slide) -> None:
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT

    r_pr = run._r.get_or_add_rPr()
    hlink = r_pr.hlinkClick
    if hlink is not None:
        rid = hlink.rId
        if rid:
            run.part.drop_rel(rid)
        r_pr._remove_hlinkClick()  # pyright: ignore[reportPrivateUsage]
    hlink = r_pr.get_or_add_hlinkClick()
    hlink.action = "ppaction://hlinksldjump"
    hlink.rId = run.part.relate_to(target_slide.part, RT.SLIDE)


def repair_executive_nav_links(prs) -> int:
    """Re-bind executive summary 'View …' links to the correct deck slides."""
    exec_slide, exec_idx = _find_executive_summary_slide(prs)
    if exec_slide is None or exec_idx is None:
        return 0

    exclude = {exec_idx}
    repaired = 0
    for pattern, keywords in _EXEC_NAV_LINKS:
        target_idx = _find_slide_by_keywords(prs, keywords, exclude_indices=exclude)
        if target_idx is None:
            logger.warning("Board PPTX: could not resolve nav target for %s", pattern.pattern)
            continue
        target_slide = prs.slides[target_idx]
        for shape in _iter_shapes(exec_slide.shapes):
            if not hasattr(shape, "text_frame") or shape.text_frame is None:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if not pattern.search(run.text or ""):
                        continue
                    _set_run_target_slide(run, target_slide)
                    repaired += 1
    if repaired:
        logger.info("Board PPTX: repaired %d executive nav hyperlinks", repaired)
    return repaired


def _replace_periods_in_shape(shape, replacements: list[tuple[re.Pattern[str], str]]) -> bool:
    if not hasattr(shape, "text_frame") or shape.text_frame is None:
        return False
    if _is_protected_nav_shape(shape):
        return False

    changed = False
    for paragraph in shape.text_frame.paragraphs:
        if _paragraph_has_slide_link(paragraph):
            continue
        raw = paragraph.text
        if not raw.strip() or _is_exec_nav_text(raw):
            continue
        updated = raw
        for pattern, repl in replacements:
            updated = pattern.sub(repl, updated)
        if updated != raw:
            paragraph.text = updated
            changed = True
    return changed


def _set_shape_text(shape, text: str) -> None:
    if not hasattr(shape, "text_frame") or shape.text_frame is None:
        return
    tf = shape.text_frame
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    if not lines:
        if tf.paragraphs:
            tf.paragraphs[0].text = ""
        else:
            shape.text = ""
        return
    for idx, line in enumerate(lines):
        if idx < len(tf.paragraphs):
            tf.paragraphs[idx].text = line
        else:
            tf.add_paragraph().text = line
    for extra in tf.paragraphs[len(lines) :]:
        extra.text = ""


def _normalize_bullet_line(text: str) -> str:
    stripped = text.strip()
    if not stripped:
        return ""
    body = stripped.lstrip("•\u2022-\u2013 \t")
    return f"• {body}" if body else ""


def _chars_per_line_for_shape(shape) -> int:
    """Rough character capacity per line from box width (10pt body text)."""
    width = max(int(getattr(shape, "width", 0) or 0), 800000)
    return max(18, int(width / 52000))


def _max_lines_for_shape(shape) -> int:
    height = max(int(getattr(shape, "height", 0) or 0), 200000)
    line_px = 127000  # ~10pt line height in EMU
    return max(1, min(3, height // line_px))


def _truncate_to_char_limit(text: str, max_chars: int) -> str:
    stripped = text.strip()
    if not stripped or max_chars <= 0:
        return ""
    if len(stripped) <= max_chars:
        return stripped
    if max_chars <= 1:
        return stripped[:max_chars]
    trimmed = stripped[: max_chars - 1].rstrip()
    if " " in trimmed:
        trimmed = trimmed.rsplit(" ", 1)[0]
    return trimmed.rstrip(".,;:") + "…"


def _truncate_text_for_shape(text: str, shape, *, max_chars: int | None = None) -> str:
    if not text.strip():
        return ""
    per_line = _chars_per_line_for_shape(shape)
    max_lines = _max_lines_for_shape(shape)
    box_cap = per_line * max_lines
    cap = max_chars if max_chars is not None else box_cap
    cap = min(cap, box_cap) if max_chars is None else min(cap, box_cap, max_chars)
    return _truncate_to_char_limit(text, cap)


def _configure_text_frame_fit(text_frame) -> None:
    """Single spacing, word wrap, and inset margins so text stays inside template boxes."""
    try:
        text_frame.word_wrap = True
    except Exception:
        pass
    for margin_attr, value in (
        ("margin_left", _TEXT_FRAME_MARGIN_EMU),
        ("margin_right", _TEXT_FRAME_MARGIN_EMU),
        ("margin_top", _TEXT_FRAME_MARGIN_EMU // 2),
        ("margin_bottom", _TEXT_FRAME_MARGIN_EMU // 2),
    ):
        try:
            setattr(text_frame, margin_attr, value)
        except Exception:
            pass
    for paragraph in text_frame.paragraphs:
        try:
            paragraph.line_spacing = 1.0
        except Exception:
            pass
        try:
            p_pr = paragraph._p.get_or_add_pPr()
            spacing = p_pr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}lnSpc")
            if spacing is None:
                from pptx.oxml.ns import qn

                spacing = p_pr.makeelement(qn("a:lnSpc"))
                p_pr.append(spacing)
            spacing.clear()
            from pptx.oxml.ns import qn

            sp_pct = spacing.makeelement(qn("a:spcPct"))
            sp_pct.set("val", "100000")
            spacing.append(sp_pct)
        except Exception:
            pass


def _estimated_shape_text_height(shape, text: str) -> int:
    if not text.strip():
        return int(getattr(shape, "height", 0) or 200000)
    per_line = _chars_per_line_for_shape(shape)
    lines = max(1, (len(text) + per_line - 1) // per_line)
    line_h = 127000
    padding = _TEXT_FRAME_MARGIN_EMU
    return lines * line_h + padding


def _reflow_key_takeaway_slots(label, slots: list[Any]) -> None:
    """Stack Key Takeaways boxes by rendered text height so bullets do not overlap."""
    if not label or not slots:
        return
    start_top = label.top + label.height + 80000
    current_top = start_top
    for shape in slots:
        text = _shape_text(shape)
        block_h = max(int(shape.height), _estimated_shape_text_height(shape, text))
        shape.top = current_top
        current_top += block_h + _KEY_TAKEAWAY_SLOT_GAP_EMU


def _set_paragraph_text_preserve_format(paragraph, text: str) -> None:
    """Replace text in-place so the template run/paragraph formatting is kept."""
    if paragraph.runs:
        paragraph.runs[0].text = text
        for run in paragraph.runs[1:]:
            run.text = ""
    else:
        paragraph.text = text


def _set_shape_single_line(shape, text: str, *, as_bullet: bool = False) -> None:
    """Write one line into a fixed template text box (single paragraph)."""
    if not hasattr(shape, "text_frame") or shape.text_frame is None:
        return
    tf = shape.text_frame
    line = _normalize_bullet_line(text) if as_bullet else text.strip()
    _configure_text_frame_fit(tf)
    if tf.paragraphs:
        _set_paragraph_text_preserve_format(tf.paragraphs[0], line)
        for extra in tf.paragraphs[1:]:
            extra.text = ""
    elif line:
        shape.text = line


def _set_fitted_single_line(
    shape,
    text: str,
    *,
    as_bullet: bool = False,
    max_chars: int | None = None,
) -> None:
    line = text.strip()
    if line:
        line = _truncate_text_for_shape(line, shape, max_chars=max_chars)
    _set_shape_single_line(shape, line, as_bullet=as_bullet)


def _commentary_text_from_slide(slide_comm) -> str:
    from app.services.reporting.export.board_api_prompts import format_key_takeaway_bullets

    if slide_comm.bullets:
        formatted = format_key_takeaway_bullets(slide_comm.bullets)
        if formatted.strip():
            return formatted.strip()
    return (slide_comm.what_happened or slide_comm.narrative_block() or "").strip()


# Semantic map: board deck slide_key -> keywords found in template slide text.
_TEMPLATE_SLIDE_KEYWORDS: dict[str, tuple[str, ...]] = {
    "executive_summary": (
        "executive summary",
        "executive operating",
        "operating summary",
        "where we stand",
    ),
    "arr_waterfall": (
        "arr analysis",
        "arr / mrr",
        "mrr waterfall",
        "arr waterfall",
        "arr roll",
        "net new arr",
    ),
    "gaap_revenue": (
        "p&l review",
        "income statement",
        "department spend",
        "gaap revenue",
        "management p&l",
    ),
    "cash_forecast": (
        "cash & liquidity",
        "cash and liquidity",
        "cash forecast",
        "liquidity",
        "cash bridge",
    ),
    "cash_flow_statement": (
        "cash flow statement",
        "operating activities",
        "investing activities",
        "financing activities",
    ),
    "gtm_performance": (
        "gtm & marketing",
        "gtm and marketing",
        "marketing efficiency",
        "channel efficiency",
        "gtm performance",
        "pipeline & channel",
    ),
    "gtm_funnel": (
        "funnel analysis",
        "new logo funnel",
        "expansion funnel",
        "pipeline to bookings",
    ),
    "headcount": (
        "workforce & headcount",
        "workforce and headcount",
        "hiring plan",
        "quota capacity",
    ),
    "risks_opportunities": (
        "risks & opportunities",
        "risks and opportunities",
        "risk / opportunity",
        "strategic assessment",
        "decision matrix",
    ),
    "financial_outlook": (
        "financial outlook",
        "h2 2026 strategy",
        "full-year forecast",
        "board discussion",
    ),
    "board_actions": (
        "board actions",
        "approvals & next steps",
        "decisions required",
    ),
}


def map_template_slides_to_slide_keys(prs) -> dict[str, int]:
    """Map board slide_key -> 1-based PPTX slide index using title/body keywords."""
    mapping: dict[str, int] = {}
    used_zero_based: set[int] = set()
    for slide_key, keywords in _TEMPLATE_SLIDE_KEYWORDS.items():
        idx0 = _find_slide_by_keywords(prs, keywords, exclude_indices=used_zero_based)
        if idx0 is None:
            continue
        mapping[slide_key] = idx0 + 1
        used_zero_based.add(idx0)
    return mapping


def _build_template_commentary_updates(
    bundle: ReportingBundle,
    prs,
    *,
    use_ai_commentary: bool,
) -> dict[int, TemplateCommentaryUpdate]:
    """Build commentary text per template slide index (API-spec path when AI enabled)."""
    from app.services.reporting.export.board_commentary_service import (
        build_slide_commentary,
        enrich_slide_with_ai,
    )

    key_to_idx = map_template_slides_to_slide_keys(prs)
    if not key_to_idx:
        logger.warning("Board PPTX: no keyword slide mapping — using outline fallback")
        outline = extract_template_outline(prs)
        return _fallback_commentary_by_slide(bundle, outline, use_ai=use_ai_commentary)

    updates: dict[int, TemplateCommentaryUpdate] = {}
    for slide_key, slide_idx in key_to_idx.items():
        base = build_slide_commentary(bundle, slide_key)
        comm = enrich_slide_with_ai(bundle, slide_key, base) if use_ai_commentary else base
        text = _commentary_text_from_slide(comm)
        bullets = tuple(comm.bullets) if comm.bullets else ()
        if text or bullets:
            updates[slide_idx] = TemplateCommentaryUpdate(
                slide_key=slide_key,
                text=text,
                bullets=bullets,
            )

    logger.info(
        "Board PPTX template commentary: %d slides (%s, ai=%s)",
        len(updates),
        ", ".join(sorted(key_to_idx.keys())),
        use_ai_commentary,
    )
    return updates


def _safe_set_shape_text(shape, text: str) -> bool:
    try:
        _set_shape_text(shape, text)
        return True
    except Exception:
        return False


def roll_template_period_labels(prs, bundle: ReportingBundle) -> int:
    replacements = _period_replacements(bundle)
    changed = 0
    for slide in prs.slides:
        for shape in _iter_shapes(slide.shapes):
            if _replace_periods_in_shape(shape, replacements):
                changed += 1
    return changed


def extract_template_outline(prs) -> list[TemplateSlideOutline]:
    outline: list[TemplateSlideOutline] = []
    for idx, slide in enumerate(prs.slides, start=1):
        texts = [_shape_text(s) for s in _iter_shapes(slide.shapes) if _shape_text(s)]
        title = texts[0][:120] if texts else f"Slide {idx}"
        preview = " | ".join(texts[:3])[:400]
        outline.append(TemplateSlideOutline(index=idx, title=title, preview=preview))
    return outline


def _fallback_commentary_by_slide(
    bundle: ReportingBundle,
    outline: list[TemplateSlideOutline],
    *,
    use_ai: bool,
) -> dict[int, TemplateCommentaryUpdate]:
    """Legacy outline-index fallback when keyword mapping finds no slides."""
    from app.services.reporting.export.board_commentary_service import build_all_slide_commentary

    commentary = build_all_slide_commentary(bundle, use_ai=use_ai)
    slide_keys = list(commentary.keys())
    updates: dict[int, TemplateCommentaryUpdate] = {}
    for i, item in enumerate(outline):
        key = slide_keys[i % len(slide_keys)] if slide_keys else "executive_summary"
        comm = commentary.get(key)
        if not comm:
            continue
        text = _commentary_text_from_slide(comm)
        bullets = tuple(comm.bullets) if comm.bullets else ()
        if text or bullets:
            updates[item.index] = TemplateCommentaryUpdate(
                slide_key=key,
                text=text,
                bullets=bullets,
            )
    return updates


def _generate_claude_template_commentary(
    bundle: ReportingBundle,
    outline: list[TemplateSlideOutline],
) -> dict[int, str]:
    settings = get_settings()
    if not (settings.anthropic_api_key or settings.openai_api_key):
        return {}

    outline_json = [
        {"index": o.index, "title": o.title, "preview": o.preview[:240]} for o in outline
    ]
    try:
        metrics_blob = ""
        try:
            metrics_blob = copilot_context_blob(bundle)[:28000]
        except Exception as exc:
            logger.warning("Template export metrics context failed: %s", exc)

        client = build_commentary_llm_client()
        raw = client.generate(
            system_prompt=(
                requirements_prompt_block("mda_deck_pptx")
                + "\n\nYou update an existing board deck template for a new close month. "
                "Keep each slide's purpose aligned to its title. Evidence-only metrics. "
                "Do not remove chart placeholders — only provide narrative commentary text."
            ),
            user_prompt=(
                f"Close: {bundle.period_label} ({bundle.as_of_period}). "
                f"Organization: {bundle.organization_name or 'SMPL'}.\n\n"
                f"Template slides:\n{json.dumps(outline_json, indent=2)}\n\n"
                f"Live metrics:\n{metrics_blob}\n\n"
                'Respond JSON only: {"slides": [{"index": 1, "commentary": "..."}, ...]} '
                "One entry per template slide index. Each commentary 2-4 sentences, board-ready."
            ),
            max_tokens=4096,
        )
        if not isinstance(raw, dict):
            return {}
        slides = raw.get("slides") or raw.get("slide_commentary") or []
        updates: dict[int, str] = {}
        for block in slides:
            if not isinstance(block, dict):
                continue
            idx = int(block.get("index") or block.get("slide") or 0)
            text = str(block.get("commentary") or block.get("narrative") or "").strip()
            if idx and text:
                updates[idx] = text
        return updates
    except Exception as exc:
        logger.warning("Template Claude commentary failed: %s", exc)
        return {}


def _is_boilerplate_shape(text: str) -> bool:
    """Footer, page number, and other shapes that must never receive commentary."""
    lowered = text.strip().lower()
    if not lowered:
        return True
    if re.search(r"\d+/\d+\s*$", lowered):
        return True
    if "confidential" in lowered and "board operating review" in lowered:
        return True
    if "not for distribution" in lowered:
        return True
    if lowered.startswith("smpl · board operating review"):
        return True
    return False


def _is_takeaway_bullet_text(text: str) -> bool:
    stripped = text.strip()
    if not stripped:
        return False
    if stripped[0] in "•\u2022":
        return True
    return stripped.startswith("- ")


def _looks_like_metric_cell(text: str) -> bool:
    stripped = text.strip()
    if not stripped:
        return True
    if re.fullmatch(r"[\$\(\)\d\.,\+%x\-bps]+", stripped.replace(" ", "")):
        return True
    if len(stripped) <= 8 and re.search(r"\d", stripped):
        return True
    return False


def _commentary_lines(text: str) -> list[str]:
    """Split AI output into lines for one shape each (bullets or sentences)."""
    lines = _split_key_takeaway_lines(text)
    if len(lines) > 1:
        return lines
    if not lines:
        return []
    single = lines[0]
    if "|" in single:
        parts = [p.strip() for p in single.split("|") if p.strip()]
        if len(parts) > 1:
            return parts
    if len(single) > 120 and ". " in single:
        sentences = [s.strip() for s in re.split(r"(?<=[.!?])\s+", single) if s.strip()]
        if len(sentences) > 1:
            return sentences
    return lines


def _split_key_takeaway_lines(text: str) -> list[str]:
    lines: list[str] = []
    for line in text.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.lower() == "key takeaways":
            continue
        lines.append(stripped)
    return lines


def _apply_lines_to_shapes(
    shapes: list[Any],
    lines: list[str],
    *,
    as_bullet: bool = False,
    max_chars: int | None = None,
) -> bool:
    if not shapes:
        return False
    for idx, shape in enumerate(shapes):
        line = lines[idx] if idx < len(lines) else ""
        _set_fitted_single_line(shape, line, as_bullet=as_bullet, max_chars=max_chars)
    return True


def _find_key_takeaways_label(slide) -> Any | None:
    for shape in _iter_shapes(slide.shapes):
        if _shape_text(shape).strip().lower() == "key takeaways":
            return shape
    return None


def _find_key_takeaways_slots(slide) -> list[Any]:
    """Fixed template text boxes under Key Takeaways — one slot per bullet, top-to-bottom."""
    label = _find_key_takeaways_label(slide)
    if label is None:
        return _find_key_takeaways_bullet_shapes(slide)

    bullets = _find_key_takeaways_bullet_shapes(slide)
    col_left = bullets[0].left if bullets else label.left
    slot_width = bullets[0].width if bullets else 2800000
    min_top = label.top + int(label.height * 0.85)
    max_top = min_top + 2800000
    tol = 350000
    width_tol = 500000
    min_width = 2000000

    candidates: list[tuple[int, int, Any]] = []
    for shape in _iter_shapes(slide.shapes):
        if shape is label:
            continue
        if not hasattr(shape, "text_frame") or shape.text_frame is None:
            continue
        if shape.top < min_top or shape.top > max_top:
            continue
        if abs(shape.left - col_left) > tol:
            continue
        if shape.width < min_width:
            continue
        if bullets and abs(shape.width - slot_width) > width_tol:
            continue
        if _is_protected_nav_shape(shape):
            continue
        text = _shape_text(shape)
        if text and _is_boilerplate_shape(text):
            continue
        if text and not _is_takeaway_bullet_text(text) and len(text) < 60:
            continue
        candidates.append((shape.top, shape.left, shape))

    if candidates:
        candidates.sort(key=lambda x: (x[0], x[1]))
        return [shape for _, _, shape in candidates]
    return _find_key_takeaways_bullet_shapes(slide)


def _find_key_takeaways_bullet_shapes(slide) -> list[Any]:
    """Return ordered text shapes that hold Key Takeaways bullets (one bullet per shape)."""
    candidates: list[tuple[int, int, Any]] = []
    for shape in _iter_shapes(slide.shapes):
        if _is_protected_nav_shape(shape):
            continue
        text = _shape_text(shape)
        if not text or _is_boilerplate_shape(text):
            continue
        if _is_takeaway_bullet_text(text):
            candidates.append((shape.top, shape.left, shape))
    candidates.sort(key=lambda item: (item[0], item[1]))
    return [shape for _, _, shape in candidates]


def _find_gtm_channel_rows(slide) -> list[tuple[int, Any]]:
    """Left-column channel name rows on the GTM performance slide."""
    rows: list[tuple[int, Any]] = []
    skip = {
        "channel",
        "commentary",
        "act spend",
        "actual spend",
        "budget spend",
        "cac",
        "pipeline",
        "marketing channel performance",
    }
    for shape in _iter_shapes(slide.shapes):
        if abs(shape.left - _GTM_CHANNEL_COL_LEFT) > 250000:
            continue
        text = _shape_text(shape)
        if not text or len(text) > 28:
            continue
        low = text.strip().lower()
        if low in skip or _looks_like_metric_cell(text):
            continue
        rows.append((shape.top, shape))
    rows.sort(key=lambda x: x[0])
    return rows


def _find_commentary_column_shapes(slide) -> list[Any]:
    """GTM channel table — right-hand Commentary column cells, row-aligned to channels."""
    header_left: int | None = None
    header_top = 0
    for shape in _iter_shapes(slide.shapes):
        if _shape_text(shape).strip().lower() == "commentary":
            header_left = shape.left
            header_top = shape.top
            break
    if header_left is None:
        return []
    tol = 250000
    cells: list[tuple[int, Any]] = []
    for shape in _iter_shapes(slide.shapes):
        if abs(shape.left - header_left) > tol:
            continue
        if shape.top <= header_top:
            continue
        if not hasattr(shape, "text_frame") or shape.text_frame is None:
            continue
        text = _shape_text(shape)
        if text and (_is_boilerplate_shape(text) or _looks_like_metric_cell(text)):
            continue
        cells.append((shape.top, shape))
    cells.sort(key=lambda x: x[0])

    channels = _find_gtm_channel_rows(slide)
    if not channels:
        return [s for _, s in cells]

    aligned: list[Any] = []
    used: set[int] = set()
    for ch_top, _ in channels:
        best_idx: int | None = None
        best_delta = _GTM_ROW_TOP_TOLERANCE + 1
        for idx, (cell_top, cell_shape) in enumerate(cells):
            if idx in used:
                continue
            delta = abs(cell_top - ch_top)
            if delta < best_delta:
                best_delta = delta
                best_idx = idx
        if best_idx is not None and best_delta <= _GTM_ROW_TOP_TOLERANCE:
            used.add(best_idx)
            aligned.append(cells[best_idx][1])
    if aligned:
        return aligned
    return [s for _, s in cells]


def _find_risk_description_shapes(slide) -> list[Any]:
    """Risks & Opportunities matrix — long description fields (not row titles)."""
    rows: list[tuple[int, int, Any]] = []
    for shape in _iter_shapes(slide.shapes):
        if _is_protected_nav_shape(shape):
            continue
        text = _shape_text(shape)
        if len(text) < 80:
            continue
        if _is_boilerplate_shape(text) or _is_takeaway_bullet_text(text):
            continue
        low = text.lower()
        if low.startswith("risk ·") or low.startswith("opp ·"):
            continue
        if shape.left < 400000:
            continue
        rows.append((shape.top, shape.left, shape))
    rows.sort(key=lambda x: (x[0], x[1]))
    return [s for _, _, s in rows]


def _find_narrative_block_shapes(slide, *, min_chars: int = 90, slide_key: str = "") -> list[Any]:
    """Funnel callouts, outlook paragraphs, board action descriptions."""
    blocks: list[tuple[int, int, Any]] = []
    for shape in _iter_shapes(slide.shapes):
        if _is_protected_nav_shape(shape):
            continue
        if not hasattr(shape, "text_frame") or shape.text_frame is None:
            continue
        text = _shape_text(shape)
        if text and len(text) < min_chars:
            continue
        if text and (_is_boilerplate_shape(text) or _is_takeaway_bullet_text(text)):
            continue
        stripped = text.strip()
        low = stripped.lower()
        if stripped.isupper() and len(stripped) < 50:
            continue
        if low.startswith("for approval") or low.startswith("for discussion"):
            continue
        if re.match(r"^\d{2}\s+for\s+(approval|discussion)", low):
            continue
        if low.startswith("owner:") or low.startswith("due:"):
            continue
        if text and _looks_like_metric_cell(stripped):
            continue
        if slide_key == "board_actions" and shape.height < 240000:
            continue
        blocks.append((shape.top, shape.left, shape))
    blocks.sort(key=lambda x: (x[0], x[1]))
    return [s for _, _, s in blocks]


def _split_narrative_into_bullets(text: str, max_count: int, char_limit: int) -> list[str]:
    """When the API returns prose instead of bullets, split into slot-sized chunks."""
    lines = [_normalize_bullet_line(line) for line in _split_key_takeaway_lines(text) if line.strip()]
    if len(lines) >= 2:
        return [_truncate_to_char_limit(line, char_limit) for line in lines[:max_count]]
    single = " ".join(line.lstrip("•").strip() for line in lines)
    if not single:
        return []
    if len(single) <= char_limit:
        return [_normalize_bullet_line(single)]
    sentences = [s.strip() for s in re.split(r"(?<=[.!?])\s+", single) if s.strip()]
    if len(sentences) >= 2:
        return [
            _truncate_to_char_limit(_normalize_bullet_line(s), char_limit)
            for s in sentences[:max_count]
        ]
    chunks: list[str] = []
    words = single.split()
    current: list[str] = []
    for word in words:
        candidate = " ".join(current + [word])
        if len(candidate) > char_limit - 1 and current:
            chunks.append(_normalize_bullet_line(" ".join(current)))
            current = [word]
        else:
            current.append(word)
    if current:
        chunks.append(_normalize_bullet_line(" ".join(current)))
    return [_truncate_to_char_limit(c, char_limit) for c in chunks[:max_count]]


def _key_takeaway_bullet_lines(update: TemplateCommentaryUpdate, slide_key: str = "") -> list[str]:
    """One string per template slot; split long prose when bullets are missing."""
    char_limit = _SLIDE_SLOT_CHAR_LIMIT.get(slide_key, 100)
    if update.bullets:
        return [
            _truncate_to_char_limit(_normalize_bullet_line(b), char_limit)
            for b in update.bullets
            if b.strip()
        ]
    raw_lines = [_normalize_bullet_line(line) for line in _split_key_takeaway_lines(update.text)]
    if len(raw_lines) == 1 and len(raw_lines[0]) > char_limit + 20:
        slots_hint = 4
        return _split_narrative_into_bullets(update.text, slots_hint, char_limit)
    return [_truncate_to_char_limit(line, char_limit) for line in raw_lines if line]


def _apply_key_takeaway_bullets(slide, bullets: list[str], *, slide_key: str = "") -> bool:
    """Fill each pre-existing template bullet box; clear unused slots; reflow by text height."""
    slots = _find_key_takeaways_slots(slide)
    if not slots:
        return False
    char_limit = _SLIDE_SLOT_CHAR_LIMIT.get(slide_key, 100)
    label = _find_key_takeaways_label(slide)
    for idx, shape in enumerate(slots):
        line = bullets[idx] if idx < len(bullets) else ""
        _set_fitted_single_line(shape, line, as_bullet=bool(line), max_chars=char_limit)
    if label:
        _reflow_key_takeaway_slots(label, slots)
    return True


def _apply_commentary_for_slide(slide, update: TemplateCommentaryUpdate) -> bool:
    zone = _SLIDE_COMMENTARY_ZONE.get(update.slide_key, "takeaway_bullets")
    char_limit = _SLIDE_SLOT_CHAR_LIMIT.get(update.slide_key)

    if zone == "takeaway_bullets":
        lines = _key_takeaway_bullet_lines(update, update.slide_key)
        if not lines:
            slots = _find_key_takeaways_slots(slide)
            for shape in slots:
                _set_fitted_single_line(shape, "")
            return bool(slots)
        return _apply_key_takeaway_bullets(slide, lines, slide_key=update.slide_key)

    lines = _commentary_lines(update.text.strip())
    if not lines and update.bullets:
        lines = [b.strip() for b in update.bullets if b.strip()]
    if not lines:
        return False

    if zone == "commentary_column":
        shapes = _find_commentary_column_shapes(slide)
        return _apply_lines_to_shapes(shapes, lines, max_chars=char_limit)
    if zone == "risk_descriptions":
        shapes = _find_risk_description_shapes(slide)
        return _apply_lines_to_shapes(shapes, lines, max_chars=char_limit)
    if zone == "narrative_blocks":
        shapes = _find_narrative_block_shapes(slide, slide_key=update.slide_key)
        return _apply_lines_to_shapes(shapes, lines, max_chars=char_limit)
    return False


def apply_template_commentary(prs, updates: dict[int, TemplateCommentaryUpdate]) -> int:
    applied = 0
    for idx, update in updates.items():
        if idx < 1 or idx > len(prs.slides):
            continue
        slide = prs.slides[idx - 1]
        if _apply_commentary_for_slide(slide, update):
            applied += 1
        else:
            logger.warning(
                "Board PPTX: no commentary zone matched slide %d (%s)",
                idx,
                update.slide_key,
            )
    return applied


def build_pptx_from_template(
    bundle: ReportingBundle,
    *,
    use_ai_commentary: bool = True,
) -> bytes | None:
    """Return filled template bytes, or None if template missing or processing fails."""
    template_path = resolve_board_pptx_template()
    if template_path is None:
        return None

    from pptx import Presentation

    try:
        prs = Presentation(str(template_path))
        roll_template_period_labels(prs, bundle)

        updates = _build_template_commentary_updates(
            bundle,
            prs,
            use_ai_commentary=use_ai_commentary,
        )

        apply_template_commentary(prs, updates)
        repair_executive_nav_links(prs)

        buf = io.BytesIO()
        prs.save(buf)
        logger.info(
            "Board PPTX from template %s (%d slides, %d commentary updates)",
            template_path.name,
            len(prs.slides),
            len(updates),
        )
        return buf.getvalue()
    except Exception as exc:
        logger.exception("Board PPTX template processing failed (%s): %s", template_path, exc)
        return None
