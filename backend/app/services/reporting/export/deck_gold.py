"""Gold deck reference — pinned PptxGenJS script + layout fingerprint for script adaptation."""

from __future__ import annotations

import json
import logging
import shutil
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from app.core.config import _BACKEND_ROOT

logger = logging.getLogger(__name__)

ARCHIVE_DIR = _BACKEND_ROOT / "tmp" / "deck-archive"
GOLD_DIR = ARCHIVE_DIR / "gold"


def gold_period_dir(period: str) -> Path:
    return GOLD_DIR / period


def gold_script_path(period: str) -> Path:
    """Pinned reference PptxGenJS script (Option 1 — adapt, do not regenerate layout)."""
    return gold_period_dir(period) / "generate_deck_reference.js"


def gold_pptx_path(period: str) -> Path:
    return gold_period_dir(period) / "mda_deck_reference.pptx"


def gold_layout_path(period: str) -> Path:
    return gold_period_dir(period) / "layout_fingerprint.json"


def archive_script_path(period: str) -> Path:
    return ARCHIVE_DIR / period / f"generate_deck_{period}.js"


def _valid_script(path: Path) -> bool:
    return path.is_file() and path.stat().st_size > 1000


def resolve_gold_script(period: str) -> Path | None:
    path = gold_script_path(period)
    return path if _valid_script(path) else None


def resolve_reference_script(period: str) -> tuple[Path | None, str]:
    """Gold script first, then latest archived script for the period (Option 1 bootstrap)."""
    gold = resolve_gold_script(period)
    if gold:
        return gold, "gold"
    archive = archive_script_path(period)
    if _valid_script(archive):
        return archive, "archive"
    return None, ""


def resolve_gold_pptx(period: str) -> Path | None:
    path = gold_pptx_path(period)
    return path if path.is_file() else None


def pin_gold_script(source: Path, period: str) -> Path:
    """Copy a known-good generate_deck.js to the gold reference path."""
    dest = gold_script_path(period)
    dest.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(source, dest)
    logger.info("Pinned gold deck script: %s", dest)
    return dest


def pin_gold_pptx(source: Path, period: str) -> Path:
    dest = gold_pptx_path(period)
    dest.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(source, dest)
    layout = extract_layout_fingerprint(dest)
    gold_layout_path(period).write_text(
        json.dumps(layout, indent=2),
        encoding="utf-8",
    )
    logger.info("Pinned gold deck pptx + layout fingerprint: %s", dest)
    return dest


def extract_layout_fingerprint(pptx_path: Path) -> dict[str, Any]:
    """Structural fingerprint from a reference PPTX (for adaptation when JS is missing)."""
    from pptx import Presentation
    from pptx.enum.chart import XL_CHART_TYPE

    prs = Presentation(str(pptx_path))
    slides: list[dict[str, Any]] = []

    for idx, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        charts: list[dict[str, Any]] = []
        tables: list[dict[str, Any]] = []

        for shape in slide.shapes:
            if shape.has_table:
                rows = shape.table.rows
                if rows:
                    headers = [rows[0].cells[c].text.strip() for c in range(len(rows[0].cells))]
                    tables.append({"headers": headers, "row_count": len(rows)})
            elif shape.has_chart:
                ch = shape.chart
                charts.append(
                    {
                        "chart_type": XL_CHART_TYPE(ch.chart_type).name,
                        "series": [s.name for s in ch.series],
                    }
                )
            elif hasattr(shape, "text") and shape.text.strip():
                line = shape.text.strip().split("\n")[0][:120]
                if line and line not in texts:
                    texts.append(line)

        slides.append(
            {
                "slide": idx,
                "text_lead": texts[:12],
                "charts": charts,
                "tables": tables,
            }
        )

    return {
        "source_pptx": pptx_path.name,
        "slide_count": len(slides),
        "slides": slides,
        "note": (
            "Match this structure when adapting the reference script. "
            "Same slide count, chart types, series names, table column headers, and section labels."
        ),
    }


def _emu_to_inches(emu: int | None) -> float:
    return round(int(emu or 0) / 914400, 3)


def extract_pptx_blueprint(pptx_path: Path) -> dict[str, Any]:
    """Detailed slide blueprint from a reference PPTX — used to bootstrap a gold PptxGenJS script."""
    from pptx import Presentation
    from pptx.enum.chart import XL_CHART_TYPE

    prs = Presentation(str(pptx_path))
    slides: list[dict[str, Any]] = []

    for idx, slide in enumerate(prs.slides, 1):
        elements: list[dict[str, Any]] = []
        for shape in slide.shapes:
            box = {
                "x": _emu_to_inches(shape.left),
                "y": _emu_to_inches(shape.top),
                "w": _emu_to_inches(shape.width),
                "h": _emu_to_inches(shape.height),
            }
            if shape.has_table:
                rows = shape.table.rows
                table_rows = [
                    [rows[r].cells[c].text.strip() for c in range(len(rows[r].cells))]
                    for r in range(len(rows))
                ]
                elements.append({"kind": "table", "box": box, "rows": table_rows})
            elif shape.has_chart:
                ch = shape.chart
                series: list[dict[str, Any]] = []
                for s in ch.series:
                    vals: list[float | str] = []
                    try:
                        vals = [round(float(v), 4) for v in s.values]
                    except (TypeError, ValueError):
                        vals = [str(v) for v in s.values]
                    series.append({"name": s.name, "values": vals})
                cats: list[str] = []
                try:
                    cats = [str(c) for c in ch.plots[0].categories]
                except (IndexError, AttributeError):
                    pass
                elements.append(
                    {
                        "kind": "chart",
                        "box": box,
                        "chart_type": XL_CHART_TYPE(ch.chart_type).name,
                        "categories": cats,
                        "series": series,
                    }
                )
            elif hasattr(shape, "text") and shape.text.strip():
                elements.append({"kind": "text", "box": box, "text": shape.text.strip()})

        elements.sort(key=lambda e: (e["box"]["y"], e["box"]["x"]))
        slides.append({"slide": idx, "elements": elements})

    return {
        "source_pptx": pptx_path.name,
        "slide_count": len(slides),
        "slides": slides,
    }


def load_gold_layout(period: str) -> dict[str, Any] | None:
    path = gold_layout_path(period)
    if path.is_file():
        return json.loads(path.read_text(encoding="utf-8"))
    pptx = resolve_gold_pptx(period)
    if pptx:
        layout = extract_layout_fingerprint(pptx)
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(json.dumps(layout, indent=2), encoding="utf-8")
        return layout
    return None


def archive_timestamp() -> str:
    return datetime.now(timezone.utc).strftime("%Y%m%dT%H%M%SZ")


def write_timestamped_archive(period: str, script: str, pptx_bytes: bytes) -> Path:
    """Immutable timestamped archive + update latest pointer files."""
    dest = ARCHIVE_DIR / period
    dest.mkdir(parents=True, exist_ok=True)
    ts = archive_timestamp()
    js_ts = dest / f"generate_deck_{period}_{ts}.js"
    pptx_ts = dest / f"mda_deck_{period}_{ts}.pptx"
    js_ts.write_text(script, encoding="utf-8")
    pptx_ts.write_bytes(pptx_bytes)
    (dest / f"generate_deck_{period}.js").write_text(script, encoding="utf-8")
    (dest / f"mda_deck_{period}.pptx").write_bytes(pptx_bytes)
    logger.info("Archived deck artifacts: %s (%s bytes)", js_ts.name, len(pptx_bytes))
    return js_ts
