"""Render a `BoardPackage` to a binary .pptx using python-pptx.

The renderer is intentionally conservative: a single layout (title + body)
re-used per slide, native python-pptx tables and column charts. No external
templates required. Charts gracefully degrade to tables when chart data is
absent or unsupported.
"""

from __future__ import annotations

import io

from app.presentation.chart_primitives.registry import apply_primitive, chart_archetype_for_slide
from app.presentation.templates.archetypes import DEPRECATED_LAYOUTS, apply_template_to_slide
from app.services.board_package.pptx_chart_archetypes import render_hints
from app.services.board_package.pptx_layout_engine import (
    CALLOUT_ROW_HEIGHT,
    FOOTNOTE_TOP,
    KPI_GAP,
    MARGIN_L,
    MAX_KPI_CARDS,
    SLIDE_SAFE_BOTTOM,
    ZONE_FOOTER_TOP,
    ZONE_KPI_BOTTOM,
    ZONE_KPI_TOP,
    ZONE_VISUAL_TOP,
    clamp_table_rows,
    compact_bullets,
    compact_callouts,
    footer_bullet_box,
    footer_zone,
    kpi_zone,
    prepare_slide_chart,
    takeaway_box,
    visual_zone,
    visual_zone_chart,
    visual_zone_table,
)
from app.services.board_package.pptx_presentation_orchestrator import prepare_package_for_render
from app.services.board_package.pptx_visual_qa import audit_presentation
from app.services.board_package.schemas import (
    BoardPackage,
    CalloutBlock,
    ChartSpec,
    KpiCard,
    SlideContent,
    TableSpec,
)


# ---------------------------------------------------------------------------
# Styling constants
# ---------------------------------------------------------------------------

# 16:9 slide dimensions in EMU (English Metric Units used by python-pptx)
SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5

TITLE_FONT_SIZE_PT = 26
SUBTITLE_FONT_SIZE_PT = 12
BODY_FONT_SIZE_PT = 11
TABLE_FONT_SIZE_PT = 10
FOOTNOTE_FONT_SIZE_PT = 9
CHART_TITLE_PT = 10
CHART_AXIS_PT = 8

ACCENT_RGB = (31, 73, 125)  # corporate navy
SUBTLE_RGB = (90, 90, 90)
GREEN_SECTION_RGB = (46, 125, 50)  # CARET-style section label
ORANGE_COMMENT_RGB = (210, 105, 30)  # marketing / narrative column header
PRESENTATION_ENGINE_VERSION = "smpl-board-v2"


def _import_pptx():
    """Lazy import so the package only loads python-pptx when rendering."""
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches, Pt

    return {
        "Presentation": Presentation,
        "CategoryChartData": CategoryChartData,
        "RGBColor": RGBColor,
        "XL_CHART_TYPE": XL_CHART_TYPE,
        "Inches": Inches,
        "Pt": Pt,
    }


def _chart_type_for(spec: ChartSpec, pptx_ns: dict) -> "pptx.enum.chart.XL_CHART_TYPE":  # type: ignore[name-defined]
    xl = pptx_ns["XL_CHART_TYPE"]
    mapping = {
        "bar": xl.BAR_CLUSTERED,
        "column": xl.COLUMN_CLUSTERED,
        # python-pptx doesn't have a true waterfall — render as a column chart.
        "waterfall": xl.COLUMN_CLUSTERED,
        "line": xl.LINE,
    }
    return mapping.get(spec.chart_type, xl.COLUMN_CLUSTERED)


# ---------------------------------------------------------------------------
# Per-element renderers
# ---------------------------------------------------------------------------


def _add_section_label(slide, label: str | None, pptx_ns: dict) -> float:
    """Green breadcrumb; returns title top offset."""
    if not label:
        return 0.3
    Inches = pptx_ns["Inches"]
    Pt = pptx_ns["Pt"]
    RGBColor = pptx_ns["RGBColor"]
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.28), Inches(12.5), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = label.upper()
    p.runs[0].font.size = Pt(11)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(*GREEN_SECTION_RGB)
    return 0.52


def _add_title(slide, title: str, subtitle: str | None, pptx_ns: dict, *, top_in: float = 0.3) -> float:
    Inches = pptx_ns["Inches"]
    Pt = pptx_ns["Pt"]
    RGBColor = pptx_ns["RGBColor"]

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(top_in), Inches(12.5), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.runs[0].font.size = Pt(TITLE_FONT_SIZE_PT)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(*ACCENT_RGB)

    sub_top = top_in + 0.72
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(sub_top), Inches(12.5), Inches(0.35))
        stf = sub_box.text_frame
        stf.word_wrap = True
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.runs[0].font.size = Pt(SUBTITLE_FONT_SIZE_PT)
        sp.runs[0].font.color.rgb = RGBColor(*SUBTLE_RGB)
        return sub_top + 0.42
    return sub_top


def _add_body_text(
    slide,
    *,
    narrative: str | None,
    bullets: list[str],
    pptx_ns: dict,
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
) -> None:
    Inches = pptx_ns["Inches"]
    Pt = pptx_ns["Pt"]
    RGBColor = pptx_ns["RGBColor"]

    if not narrative and not bullets:
        return

    box = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )
    tf = box.text_frame
    tf.word_wrap = True

    first = True
    if narrative:
        p = tf.paragraphs[0]
        p.text = narrative
        p.runs[0].font.size = Pt(BODY_FONT_SIZE_PT)
        p.runs[0].font.color.rgb = RGBColor(20, 20, 20)
        first = False

    for line in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = f"• {line}" if line else ""
        p.level = 0
        if p.runs:
            p.runs[0].font.size = Pt(BODY_FONT_SIZE_PT)


def _add_table(
    slide,
    spec: TableSpec,
    pptx_ns: dict,
    *,
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
) -> None:
    Inches = pptx_ns["Inches"]
    Pt = pptx_ns["Pt"]
    RGBColor = pptx_ns["RGBColor"]

    rows = len(spec.rows) + 1
    cols = len(spec.headers)
    if rows == 0 or cols == 0:
        return

    shape = slide.shapes.add_table(
        rows, cols, Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )
    table = shape.table

    for c, header in enumerate(spec.headers):
        cell = table.cell(0, c)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(TABLE_FONT_SIZE_PT)
                run.font.color.rgb = RGBColor(255, 255, 255)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(*ACCENT_RGB)

    for r, row in enumerate(spec.rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(TABLE_FONT_SIZE_PT)


def _add_chart(
    slide,
    spec: ChartSpec,
    pptx_ns: dict,
    *,
    left_in: float,
    top_in: float,
    width_in: float,
    height_in: float,
    slide_id: str | None = None,
) -> None:
    Inches = pptx_ns["Inches"]
    Pt = pptx_ns["Pt"]
    CategoryChartData = pptx_ns["CategoryChartData"]

    if slide_id:
        spec = apply_primitive(spec, slide_id)
    prepared = prepare_slide_chart(spec)
    if not prepared or not prepared.categories or not prepared.series:
        return

    arch = prepared.archetype or chart_archetype_for_slide(slide_id or "", prepared)
    hints = render_hints(arch)  # type: ignore[arg-type]

    data = CategoryChartData()
    data.categories = prepared.categories
    for series_name, values in prepared.series.items():
        padded = list(values) + [0.0] * (len(prepared.categories) - len(values))
        data.add_series(series_name, padded[: len(prepared.categories)])

    chart_type = _chart_type_for(prepared, pptx_ns)
    chart = slide.shapes.add_chart(
        chart_type,
        Inches(left_in),
        Inches(top_in),
        Inches(width_in),
        Inches(height_in),
        data,
    ).chart
    chart.has_title = bool(prepared.title)
    if chart.has_title:
        chart.chart_title.text_frame.text = prepared.title[: hints.title_max_len]
        try:
            chart.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(CHART_TITLE_PT)
        except Exception:
            pass
    chart.has_legend = hints.legend and len(prepared.series) > 1
    try:
        n_cat = len(prepared.categories)
        chart.category_axis.tick_labels.font.size = Pt(7 if n_cat > 6 else CHART_AXIS_PT)
        if n_cat > 5:
            chart.category_axis.tick_label_position = 2  # low
        if prepared.y_axis_label:
            chart.value_axis.has_title = True
            chart.value_axis.axis_title.text_frame.text = prepared.y_axis_label[:20]
            chart.value_axis.axis_title.text_frame.paragraphs[0].runs[0].font.size = Pt(8)
        chart.value_axis.tick_labels.number_format = hints.number_format
        chart.value_axis.tick_labels.font.size = Pt(CHART_AXIS_PT)
        chart.value_axis.major_gridlines.format.line.fill.background()
    except Exception:
        pass


def _tone_rgb(tone: str | None) -> tuple[int, int, int]:
    return {
        "favorable": (34, 139, 34),
        "unfavorable": (180, 50, 50),
        "watch": (200, 140, 0),
        "neutral": ACCENT_RGB,
    }.get(tone or "neutral", ACCENT_RGB)


def _add_single_kpi_box(slide, card: KpiCard, pptx_ns: dict, *, left: float, top: float, width: float, height: float) -> None:
    Inches = pptx_ns["Inches"]
    Pt = pptx_ns["Pt"]
    RGBColor = pptx_ns["RGBColor"]
    dir_glyph = {"up": "▲", "down": "▼", "flat": "●"}.get(card.direction or "flat", "")

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    group = f" · {card.group.title()}" if card.group else ""
    p0.text = f"{card.label}{group}"
    p0.runs[0].font.size = Pt(9)
    p0.runs[0].font.color.rgb = RGBColor(*SUBTLE_RGB)
    p1 = tf.add_paragraph()
    p1.text = f"{dir_glyph} {card.value}".strip()
    p1.runs[0].font.size = Pt(15)
    p1.runs[0].font.bold = True
    p1.runs[0].font.color.rgb = RGBColor(*_tone_rgb(card.tone))
    if card.subtext:
        p2 = tf.add_paragraph()
        p2.text = card.subtext[:48]
        p2.runs[0].font.size = Pt(8)
        p2.runs[0].font.color.rgb = RGBColor(*SUBTLE_RGB)


def _add_kpi_cards(slide, cards: list[KpiCard], pptx_ns: dict, *, grouped: bool = False) -> float:
    """Uniform KPI row in top zone; max 4 cards."""
    if not cards:
        return ZONE_KPI_TOP + 0.9
    display = cards[:MAX_KPI_CARDS]
    n = len(display)
    kz = kpi_zone()
    gap = KPI_GAP
    card_w = (kz.width - gap * (n - 1)) / n
    left = kz.left
    top = kz.top
    height = kz.height
    for card in display:
        _add_single_kpi_box(slide, card, pptx_ns, left=left, top=top, width=card_w, height=height)
        left += card_w + gap
    return ZONE_KPI_TOP + height + 0.08


def _add_callouts(slide, callouts: list[CalloutBlock], pptx_ns: dict, *, top_in: float) -> None:
    if not callouts:
        return
    Inches = pptx_ns["Inches"]
    Pt = pptx_ns["Pt"]
    RGBColor = pptx_ns["RGBColor"]
    kind_color = {"win": (34, 139, 34), "risk": (180, 50, 50), "action": (31, 73, 125)}
    fz = footer_zone()
    n = min(len(callouts), 2)
    width = (fz.width - KPI_GAP) / n
    left = fz.left
    for i, c in enumerate(callouts[:2]):
        box = slide.shapes.add_textbox(
            Inches(left + i * (width + KPI_GAP)),
            Inches(top_in),
            Inches(width),
            Inches(CALLOUT_ROW_HEIGHT),
        )
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        prefix = {"win": "WIN", "risk": "RISK", "action": "ACTION"}.get(c.kind, c.kind.upper())
        owner = f" ({c.owner})" if c.owner else ""
        p.text = f"{prefix}{owner}: {c.text[:72]}"
        p.runs[0].font.size = Pt(8)
        p.runs[0].font.color.rgb = RGBColor(*kind_color.get(c.kind, SUBTLE_RGB))


def _add_spotlights(slide, cards: list[KpiCard], pptx_ns: dict) -> None:
    if not cards:
        return
    Inches = pptx_ns["Inches"]
    n = min(len(cards), 3)
    gap = 0.2
    w = (12.3 - gap * (n - 1)) / n
    left = 0.5
    for card in cards[:3]:
        _add_single_kpi_box(slide, card, pptx_ns, left=left, top=2.0, width=w, height=1.2)
        left += w + gap


def _add_footnote(slide, footnote: str, pptx_ns: dict) -> None:
    Inches = pptx_ns["Inches"]
    Pt = pptx_ns["Pt"]
    RGBColor = pptx_ns["RGBColor"]

    box = slide.shapes.add_textbox(Inches(0.6), Inches(FOOTNOTE_TOP), Inches(12.1), Inches(0.35))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = footnote
    p.runs[0].font.size = Pt(FOOTNOTE_FONT_SIZE_PT)
    p.runs[0].font.italic = True
    p.runs[0].font.color.rgb = RGBColor(*SUBTLE_RGB)


# ---------------------------------------------------------------------------
# Slide layout
# ---------------------------------------------------------------------------


def _add_engine_watermark(slide, pptx_ns: dict) -> None:
    """Visible proof the smpl-board-v2 renderer ran (bottom-right)."""
    Inches = pptx_ns["Inches"]
    Pt = pptx_ns["Pt"]
    RGBColor = pptx_ns["RGBColor"]
    box = slide.shapes.add_textbox(Inches(10.85), Inches(6.95), Inches(2.2), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.text = f"SMPL · {PRESENTATION_ENGINE_VERSION}"
    p.runs[0].font.size = Pt(7)
    p.runs[0].font.color.rgb = RGBColor(120, 120, 120)


def _render_cover_slide(prs, package: BoardPackage, pptx_ns: dict) -> None:
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    Inches = pptx_ns["Inches"]
    Pt = pptx_ns["Pt"]
    RGBColor = pptx_ns["RGBColor"]

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.5), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"{package.organization_name} — Board Operating Review"
    p.runs[0].font.size = Pt(40)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(*ACCENT_RGB)

    sub = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.5), Inches(0.7))
    stf = sub.text_frame
    sp = stf.paragraphs[0]
    from app.services.reporting.export.company_context import COMPANY_NAME, COMPANY_TAGLINE

    if package.organization_name.strip().upper() == COMPANY_NAME:
        sp.text = f"{COMPANY_TAGLINE}"
        sp.runs[0].font.size = Pt(16)
        sp.runs[0].font.color.rgb = RGBColor(*SUBTLE_RGB)
        sp2 = stf.add_paragraph()
        sp2.text = (
            f"{package.period_label}  |  Prepared for {package.prepared_for}  |  "
            f"{package.prepared_date.isoformat()}"
        )
        sp2.runs[0].font.size = Pt(14)
        sp2.runs[0].font.color.rgb = RGBColor(*SUBTLE_RGB)
        sp3 = stf.add_paragraph()
        sp3.text = f"Engine {PRESENTATION_ENGINE_VERSION}"
        sp3.runs[0].font.size = Pt(9)
        sp3.runs[0].font.color.rgb = RGBColor(*SUBTLE_RGB)
    else:
        sp.text = (
            f"{package.period_label}  |  Prepared for {package.prepared_for}  |  "
            f"{package.prepared_date.isoformat()}"
        )
        sp.runs[0].font.size = Pt(18)
        sp.runs[0].font.color.rgb = RGBColor(*SUBTLE_RGB)
        sp2 = stf.add_paragraph()
        sp2.text = f"Engine {PRESENTATION_ENGINE_VERSION}"
        sp2.runs[0].font.size = Pt(10)
        sp2.runs[0].font.color.rgb = RGBColor(*GREEN_SECTION_RGB)

    _add_engine_watermark(slide, pptx_ns)


def _render_section_divider(prs, slide_content: SlideContent, pptx_ns: dict) -> None:
    """Minimal section transition — reset context, one KPI, narrative hook."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    Inches = pptx_ns["Inches"]
    Pt = pptx_ns["Pt"]
    RGBColor = pptx_ns["RGBColor"]

    if slide_content.section_label:
        box = slide.shapes.add_textbox(Inches(MARGIN_L), Inches(0.55), Inches(11.5), Inches(0.24))
        p = box.text_frame.paragraphs[0]
        p.text = slide_content.section_label.upper()
        p.runs[0].font.size = Pt(10)
        p.runs[0].font.bold = True
        p.runs[0].font.color.rgb = RGBColor(*GREEN_SECTION_RGB)

    title_box = slide.shapes.add_textbox(Inches(MARGIN_L), Inches(1.55), Inches(11.5), Inches(0.75))
    p = title_box.text_frame.paragraphs[0]
    p.text = slide_content.title
    p.runs[0].font.size = Pt(30)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(*ACCENT_RGB)

    y = 2.45
    if slide_content.subtitle:
        sub = slide.shapes.add_textbox(Inches(MARGIN_L), Inches(y), Inches(11.5), Inches(0.3))
        sp = sub.text_frame.paragraphs[0]
        sp.text = slide_content.subtitle
        sp.runs[0].font.size = Pt(12)
        sp.runs[0].font.color.rgb = RGBColor(*SUBTLE_RGB)
        y += 0.42

    if slide_content.kpi_cards:
        _add_kpi_cards(slide, slide_content.kpi_cards[:1], pptx_ns)
        y = ZONE_KPI_BOTTOM + 0.35

    narrative = slide_content.key_takeaway or slide_content.narrative
    if narrative:
        nb = slide.shapes.add_textbox(Inches(MARGIN_L + 0.15), Inches(y), Inches(10.5), Inches(1.05))
        np = nb.text_frame.paragraphs[0]
        np.text = narrative[:200]
        np.runs[0].font.size = Pt(13)
        np.runs[0].font.italic = True
        np.runs[0].font.color.rgb = RGBColor(55, 55, 55)


def _add_key_takeaway(slide, text: str | None, pptx_ns: dict, *, top_in: float | None = None) -> None:
    if not text:
        return
    Inches = pptx_ns["Inches"]
    Pt = pptx_ns["Pt"]
    RGBColor = pptx_ns["RGBColor"]
    tb = takeaway_box()
    top = top_in if top_in is not None else tb.top
    box = slide.shapes.add_textbox(Inches(tb.left), Inches(top), Inches(tb.width), Inches(tb.height))
    p = box.text_frame.paragraphs[0]
    from app.services.reporting.export.board_chart_density import truncate_bullet

    p.text = f"Takeaway: {truncate_bullet(text, max_len=140)}"
    p.runs[0].font.size = Pt(10)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = RGBColor(*ACCENT_RGB)


def _render_footer_band(slide, slide_content: SlideContent, pptx_ns: dict) -> None:
    """Footer zone only — callouts, bullets, takeaway stacked without overlap."""
    bullets = compact_bullets(slide_content)
    callouts = compact_callouts(slide_content.callouts)
    fz = footer_zone()
    top = fz.top
    if callouts:
        _add_callouts(slide, callouts, pptx_ns, top_in=top)
    if bullets:
        bb = footer_bullet_box(has_callouts=bool(callouts))
        _add_body_text(
            slide,
            narrative=None,
            bullets=bullets,
            pptx_ns=pptx_ns,
            left_in=bb.left,
            top_in=bb.top,
            width_in=bb.width,
            height_in=bb.height,
        )
    if slide_content.key_takeaway:
        _add_key_takeaway(slide, slide_content.key_takeaway, pptx_ns)


def _render_executive_dashboard(slide, slide_content: SlideContent, pptx_ns: dict, body_top: float) -> None:
    """IMG 8039 — KPI table left, key takeaways right."""
    if slide_content.table and slide_content.table.rows:
        _add_table(slide, slide_content.table, pptx_ns, left_in=0.5, top_in=body_top, width_in=7.0, height_in=5.5)
    _add_body_text(
        slide,
        narrative=slide_content.narrative,
        bullets=slide_content.bullets[:7],
        pptx_ns=pptx_ns,
        left_in=7.8,
        top_in=body_top,
        width_in=5.0,
        height_in=5.5,
    )
    if slide_content.footnote:
        _add_footnote(slide, slide_content.footnote, pptx_ns)


def _add_commentary_heading(
    slide, heading: str | None, pptx_ns: dict, *, left_in: float, top_in: float, width_in: float
) -> float:
    if not heading:
        return top_in
    Inches = pptx_ns["Inches"]
    Pt = pptx_ns["Pt"]
    RGBColor = pptx_ns["RGBColor"]
    hb = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(0.28))
    hp = hb.text_frame.paragraphs[0]
    hp.text = heading
    hp.runs[0].font.size = Pt(12)
    hp.runs[0].font.bold = True
    hp.runs[0].font.color.rgb = RGBColor(*ORANGE_COMMENT_RGB)
    return top_in + 0.32


def _render_executive_scorecard(slide, slide_content: SlideContent, pptx_ns: dict, body_top: float) -> None:
    """IMG 8039 pattern: KPI row, YTD scorecard left, key takeaways right, optional ARR trend below."""
    Inches = pptx_ns["Inches"]
    _add_kpi_cards(slide, slide_content.kpi_cards, pptx_ns)
    top = ZONE_KPI_BOTTOM + 0.12
    left_w = 5.85
    right_left = MARGIN_L + left_w + 0.2
    right_w = 12.05 - right_left - MARGIN_L

    if slide_content.table and slide_content.table.rows:
        tbl = clamp_table_rows(slide_content.table, max_rows=6)
        if tbl:
            _add_table(
                slide,
                tbl,
                pptx_ns,
                left_in=MARGIN_L,
                top_in=top,
                width_in=left_w,
                height_in=2.35,
            )

    rtop = top
    rtop = _add_commentary_heading(
        slide,
        slide_content.commentary_heading or "Key takeaways",
        pptx_ns,
        left_in=right_left,
        top_in=rtop,
        width_in=right_w,
    )
    _add_body_text(
        slide,
        narrative=None,
        bullets=compact_bullets(slide_content)[:5],
        pptx_ns=pptx_ns,
        left_in=right_left,
        top_in=rtop,
        width_in=right_w,
        height_in=2.2,
    )

    chart_top = top + 2.45
    if slide_content.chart:
        _add_chart(
            slide,
            slide_content.chart,
            pptx_ns,
            left_in=MARGIN_L,
            top_in=chart_top,
            width_in=12.05,
            height_in=1.55,
            slide_id=slide_content.slide_id,
        )
    elif slide_content.key_takeaway:
        _add_key_takeaway(slide, slide_content.key_takeaway, pptx_ns, top_in=chart_top)

    if slide_content.footnote:
        _add_footnote(slide, slide_content.footnote, pptx_ns)


def _render_executive_ytd(slide, slide_content: SlideContent, pptx_ns: dict, body_top: float) -> None:
    """CFO: KPI row + single YTD trajectory chart (no stacked table below chart)."""
    _add_kpi_cards(slide, slide_content.kpi_cards, pptx_ns)
    if slide_content.chart:
        cbox = visual_zone_chart(with_table=False)
        _add_chart(
            slide,
            slide_content.chart,
            pptx_ns,
            left_in=cbox.left,
            top_in=cbox.top,
            width_in=cbox.width,
            height_in=cbox.height,
            slide_id=slide_content.slide_id,
        )
    elif slide_content.table and slide_content.table.rows:
        tbox = visual_zone_table()
        tbl = clamp_table_rows(slide_content.table, max_rows=4)
        if tbl:
            _add_table(slide, tbl, pptx_ns, left_in=tbox.left, top_in=tbox.top, width_in=tbox.width, height_in=tbox.height)
    _render_footer_band(slide, slide_content, pptx_ns)
    if slide_content.footnote:
        _add_footnote(slide, slide_content.footnote, pptx_ns)


def _render_story_slide(slide, slide_content: SlideContent, pptx_ns: dict, body_top: float) -> None:
    """One primary visual in middle zone; KPIs top; commentary footer only."""
    _add_kpi_cards(slide, slide_content.kpi_cards, pptx_ns)
    if slide_content.key_takeaway:
        Inches = pptx_ns["Inches"]
        Pt = pptx_ns["Pt"]
        RGBColor = pptx_ns["RGBColor"]
        from app.services.reporting.export.board_chart_density import truncate_bullet

        ib = slide.shapes.add_textbox(Inches(MARGIN_L), Inches(ZONE_KPI_BOTTOM + 0.02), Inches(12.0), Inches(0.3))
        ip = ib.text_frame.paragraphs[0]
        ip.text = truncate_bullet(slide_content.key_takeaway, max_len=130)
        ip.runs[0].font.size = Pt(10)
        ip.runs[0].font.italic = True
        ip.runs[0].font.color.rgb = RGBColor(*ACCENT_RGB)
    if slide_content.chart:
        cbox = visual_zone_chart(with_table=False)
        _add_chart(
            slide,
            slide_content.chart,
            pptx_ns,
            left_in=cbox.left,
            top_in=cbox.top,
            width_in=cbox.width,
            height_in=cbox.height,
            slide_id=slide_content.slide_id,
        )
    elif slide_content.table and slide_content.table.rows:
        tbox = visual_zone_table()
        tbl = clamp_table_rows(slide_content.table)
        if tbl:
            _add_table(slide, tbl, pptx_ns, left_in=tbox.left, top_in=tbox.top, width_in=tbox.width, height_in=tbox.height)
    _render_footer_band(slide, slide_content, pptx_ns)
    if slide_content.footnote:
        _add_footnote(slide, slide_content.footnote, pptx_ns)


def _render_mda_narrative(slide, slide_content: SlideContent, pptx_ns: dict, body_top: float) -> None:
    """MD&A: narrative in middle zone, footer for takeaway."""
    Inches = pptx_ns["Inches"]
    Pt = pptx_ns["Pt"]
    RGBColor = pptx_ns["RGBColor"]
    top = ZONE_VISUAL_TOP
    if slide_content.commentary_heading:
        hb = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.0), Inches(0.28))
        hp = hb.text_frame.paragraphs[0]
        hp.text = slide_content.commentary_heading
        hp.runs[0].font.size = Pt(12)
        hp.runs[0].font.bold = True
        hp.runs[0].font.color.rgb = RGBColor(*ORANGE_COMMENT_RGB)
        top += 0.35
    vbox = visual_zone()
    _add_body_text(
        slide,
        narrative=slide_content.narrative,
        bullets=compact_bullets(slide_content)[:5],
        pptx_ns=pptx_ns,
        left_in=vbox.left,
        top_in=top,
        width_in=vbox.width,
        height_in=vbox.height,
    )
    _render_footer_band(slide, slide_content, pptx_ns)


def _render_narrative_table_split(slide, slide_content: SlideContent, pptx_ns: dict, body_top: float) -> None:
    """Slide 5 / MD&A — narrative left, table/chart right."""
    Inches = pptx_ns["Inches"]
    Pt = pptx_ns["Pt"]
    RGBColor = pptx_ns["RGBColor"]
    left_w = 4.6
    right_left = 5.3
    right_w = 7.5
    if slide_content.commentary_heading:
        hb = slide.shapes.add_textbox(Inches(0.5), Inches(body_top), Inches(left_w), Inches(0.3))
        hp = hb.text_frame.paragraphs[0]
        hp.text = slide_content.commentary_heading
        hp.runs[0].font.size = Pt(13)
        hp.runs[0].font.bold = True
        hp.runs[0].font.color.rgb = RGBColor(*ORANGE_COMMENT_RGB)
        body_top += 0.35
    _add_body_text(
        slide,
        narrative=slide_content.narrative,
        bullets=slide_content.bullets[:6],
        pptx_ns=pptx_ns,
        left_in=0.5,
        top_in=body_top,
        width_in=left_w,
        height_in=5.8,
    )
    rtop = body_top
    if slide_content.chart:
        _add_chart(
            slide,
            slide_content.chart,
            pptx_ns,
            left_in=right_left,
            top_in=rtop,
            width_in=right_w,
            height_in=2.5,
            slide_id=slide_content.slide_id,
        )
        rtop += 2.65
    if slide_content.secondary_chart:
        _add_chart(
            slide,
            slide_content.secondary_chart,
            pptx_ns,
            left_in=right_left,
            top_in=rtop,
            width_in=right_w,
            height_in=2.2,
            slide_id=slide_content.slide_id,
        )
        rtop += 2.35
    if slide_content.table and slide_content.table.rows and not slide_content.chart:
        th = min(5.5, 6.8 - rtop)
        _add_table(slide, slide_content.table, pptx_ns, left_in=right_left, top_in=rtop, width_in=right_w, height_in=th)
    elif slide_content.table and slide_content.table.rows:
        th = min(2.8, max(1.5, 6.8 - rtop))
        _add_table(slide, slide_content.table, pptx_ns, left_in=right_left, top_in=rtop, width_in=right_w, height_in=th)


def _render_marketing_source(slide, slide_content: SlideContent, pptx_ns: dict, body_top: float) -> None:
    """Slide 13 — table left (~65%), orange commentary right."""
    if slide_content.table and slide_content.table.rows:
        _add_table(slide, slide_content.table, pptx_ns, left_in=0.5, top_in=body_top, width_in=7.8, height_in=5.6)
    Inches = pptx_ns["Inches"]
    Pt = pptx_ns["Pt"]
    RGBColor = pptx_ns["RGBColor"]
    cx = 8.5
    if slide_content.commentary_heading:
        hb = slide.shapes.add_textbox(Inches(cx), Inches(body_top), Inches(4.3), Inches(0.3))
        hp = hb.text_frame.paragraphs[0]
        hp.text = slide_content.commentary_heading
        hp.runs[0].font.size = Pt(13)
        hp.runs[0].font.bold = True
        hp.runs[0].font.color.rgb = RGBColor(*ORANGE_COMMENT_RGB)
        body_top += 0.35
    _add_body_text(
        slide,
        narrative=slide_content.narrative,
        bullets=slide_content.bullets[:6],
        pptx_ns=pptx_ns,
        left_in=cx,
        top_in=body_top,
        width_in=4.3,
        height_in=5.5,
    )


def _render_cash_trend(slide, slide_content: SlideContent, pptx_ns: dict, body_top: float) -> None:
    """Cash line chart only in visual zone; drivers in footer."""
    _add_kpi_cards(slide, slide_content.kpi_cards, pptx_ns)
    if slide_content.chart:
        cbox = visual_zone_chart(with_table=False)
        _add_chart(
            slide,
            slide_content.chart,
            pptx_ns,
            left_in=cbox.left,
            top_in=cbox.top,
            width_in=cbox.width,
            height_in=cbox.height,
            slide_id=slide_content.slide_id,
        )
    _render_footer_band(slide, slide_content, pptx_ns)


def _render_content_slide(prs, slide_content: SlideContent, pptx_ns: dict) -> None:
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    slide_content = apply_template_to_slide(slide_content)
    slide_layout = slide_content.layout or "story_slide"
    if slide_layout in DEPRECATED_LAYOUTS:
        slide_layout = DEPRECATED_LAYOUTS[slide_layout]

    if slide_layout in ("section_divider", "section_transition"):
        _render_section_divider(prs, slide_content, pptx_ns)
        return

    title_top = _add_section_label(slide, slide_content.section_label, pptx_ns)
    body_top = _add_title(slide, slide_content.title, slide_content.subtitle, pptx_ns, top_in=title_top)

    if slide_layout == "executive_dashboard":
        _render_executive_dashboard(slide, slide_content, pptx_ns, body_top)
        if slide_content.footnote:
            _add_footnote(slide, slide_content.footnote, pptx_ns)
        return

    if slide_layout == "executive_scorecard":
        _render_executive_scorecard(slide, slide_content, pptx_ns, body_top)
        return

    if slide_layout == "executive_ytd":
        _render_executive_ytd(slide, slide_content, pptx_ns, body_top)
        return

    if slide_layout == "story_slide":
        _render_story_slide(slide, slide_content, pptx_ns, body_top)
        return

    if slide_layout == "mda_narrative":
        _render_mda_narrative(slide, slide_content, pptx_ns, body_top)
        return

    if slide_layout == "narrative_table_split":
        _render_narrative_table_split(slide, slide_content, pptx_ns, body_top)
        if slide_content.footnote:
            _add_footnote(slide, slide_content.footnote, pptx_ns)
        return

    if slide_layout == "marketing_source":
        _render_marketing_source(slide, slide_content, pptx_ns, body_top)
        if slide_content.footnote:
            _add_footnote(slide, slide_content.footnote, pptx_ns)
        return

    if slide_layout == "cash_trend":
        _render_cash_trend(slide, slide_content, pptx_ns, body_top)
        if slide_content.footnote:
            _add_footnote(slide, slide_content.footnote, pptx_ns)
        return

    if slide_layout == "spotlight":
        vbox = visual_zone_chart(with_table=False)
        n = min(len(slide_content.spotlight_cards), 3)
        gap = 0.15
        w = (vbox.width - gap * (n - 1)) / max(n, 1)
        left = vbox.left
        for card in slide_content.spotlight_cards[:3]:
            _add_single_kpi_box(slide, card, pptx_ns, left=left, top=vbox.top, width=w, height=min(2.0, vbox.height))
            left += w + gap
        _render_footer_band(slide, slide_content, pptx_ns)
        if slide_content.footnote:
            _add_footnote(slide, slide_content.footnote, pptx_ns)
        return

    if slide_layout == "risk_matrix":
        _add_callouts(slide, compact_callouts(slide_content.callouts), pptx_ns, top_in=ZONE_VISUAL_TOP)
        _render_footer_band(slide, slide_content, pptx_ns)
        if slide_content.footnote:
            _add_footnote(slide, slide_content.footnote, pptx_ns)
        return

    has_kpis = bool(slide_content.kpi_cards)
    body_top = 1.6
    if has_kpis:
        body_top = _add_kpi_cards(slide, slide_content.kpi_cards, pptx_ns, grouped=False)

    has_table = slide_content.table is not None and slide_content.table.rows
    has_chart = slide_content.chart is not None
    has_secondary = slide_content.secondary_chart is not None

    body_left = 0.5
    body_height = max(2.8, 6.8 - body_top)

    if slide_layout == "dual_metric" and has_chart:
        _add_chart(
            slide,
            slide_content.chart,
            pptx_ns,
            left_in=0.5,
            top_in=body_top,
            width_in=6.0,
            height_in=2.4,
            slide_id=slide_content.slide_id,
        )
        if has_secondary:
            _add_chart(
                slide,
                slide_content.secondary_chart,
                pptx_ns,
                left_in=6.7,
                top_in=body_top,
                width_in=6.0,
                height_in=2.4,
                slide_id=slide_content.slide_id,
            )
        _add_body_text(
            slide,
            narrative=slide_content.narrative,
            bullets=slide_content.bullets[:5],
            pptx_ns=pptx_ns,
            left_in=0.5,
            top_in=body_top + 2.55,
            width_in=12.3,
            height_in=6.8 - (body_top + 2.55),
        )
        if slide_content.footnote:
            _add_footnote(slide, slide_content.footnote, pptx_ns)
        return

    if slide_layout in ("compact_table", "narrative") and not has_chart:
        if has_table:
            _add_table(slide, slide_content.table, pptx_ns, left_in=6.4, top_in=body_top, width_in=6.4, height_in=body_height)
        _add_body_text(
            slide,
            narrative=slide_content.narrative,
            bullets=slide_content.bullets[:6],
            pptx_ns=pptx_ns,
            left_in=0.5,
            top_in=body_top,
            width_in=5.8 if has_table else 12.3,
            height_in=body_height,
        )
        if slide_content.footnote:
            _add_footnote(slide, slide_content.footnote, pptx_ns)
        return

    if has_table and has_chart:
        # Left column: narrative + bullets. Right column: chart on top, table below.
        _add_body_text(
            slide,
            narrative=slide_content.narrative,
            bullets=slide_content.bullets,
            pptx_ns=pptx_ns,
            left_in=body_left,
            top_in=body_top,
            width_in=5.5,
            height_in=body_height,
        )
        _add_chart(
            slide,
            slide_content.chart,
            pptx_ns,
            left_in=6.4,
            top_in=body_top,
            width_in=6.4,
            height_in=2.8,
            slide_id=slide_content.slide_id,
        )
        _add_table(
            slide,
            slide_content.table,
            pptx_ns,
            left_in=6.4,
            top_in=body_top + 3.0,
            width_in=6.4,
            height_in=2.2,
        )
    elif has_chart:
        _add_body_text(
            slide,
            narrative=slide_content.narrative,
            bullets=slide_content.bullets,
            pptx_ns=pptx_ns,
            left_in=body_left,
            top_in=body_top,
            width_in=5.5,
            height_in=body_height,
        )
        _add_chart(
            slide,
            slide_content.chart,
            pptx_ns,
            left_in=6.4,
            top_in=body_top,
            width_in=6.4,
            height_in=body_height,
            slide_id=slide_content.slide_id,
        )
    elif has_table:
        _add_body_text(
            slide,
            narrative=slide_content.narrative,
            bullets=slide_content.bullets,
            pptx_ns=pptx_ns,
            left_in=body_left,
            top_in=body_top,
            width_in=5.5,
            height_in=body_height,
        )
        _add_table(
            slide,
            slide_content.table,
            pptx_ns,
            left_in=6.4,
            top_in=body_top,
            width_in=6.4,
            height_in=body_height,
        )
    else:
        _add_body_text(
            slide,
            narrative=slide_content.narrative,
            bullets=slide_content.bullets,
            pptx_ns=pptx_ns,
            left_in=body_left,
            top_in=body_top,
            width_in=12.3,
            height_in=body_height,
        )

    if slide_content.footnote:
        _add_footnote(slide, slide_content.footnote, pptx_ns)


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------


def render_pptx_bytes(package: BoardPackage, *, skip_orchestration: bool = False) -> bytes:
    """Build a PowerPoint presentation in memory and return its bytes."""
    if not skip_orchestration:
        package = prepare_package_for_render(package)

    pptx_ns = _import_pptx()
    Inches = pptx_ns["Inches"]

    prs = pptx_ns["Presentation"]()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)

    _render_cover_slide(prs, package, pptx_ns)
    for slide_content in package.slides:
        _render_content_slide(prs, slide_content, pptx_ns)
        if prs.slides:
            _add_engine_watermark(prs.slides[-1], pptx_ns)

    audit_presentation(prs)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
